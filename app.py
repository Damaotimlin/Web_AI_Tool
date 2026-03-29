import gradio as gr
import requests
import http.cookiejar
import browser_cookie3
import json
import subprocess
import os
import re
import threading
from urllib.parse import urljoin, urlparse
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import datetime
import prompts
import importlib


def get_prompts():
    """Hot-reload prompts.py so edits take effect without restarting the server."""
    importlib.reload(prompts)
    return prompts


class StopFlag:
    """Lightweight cancellation flag for long-running generators."""
    def __init__(self):
        self._event = threading.Event()

    def stop(self):
        self._event.set()

    @property
    def stopped(self):
        return self._event.is_set()

    def check(self):
        """Raise GeneratorExit if stopped — call this between steps."""
        if self._event.is_set():
            raise GeneratorExit("Stopped by user")


class SharedProgress:
    """Thread-safe shared progress visible to all connected clients via polling."""
    def __init__(self):
        self._lock = threading.Lock()
        self._state = {
            "research_log": "", "research_kw": "", "research_out": "",
            "pick_log": "", "keynote_log": "",
            "running": None,  # "research" | "pick" | "keynote" | None
        }
        self._version = 0  # bumped on every update so poll can detect changes

    def update(self, **kwargs):
        with self._lock:
            self._state.update(kwargs)
            self._version += 1

    def get(self, key, default=""):
        with self._lock:
            return self._state.get(key, default)

    def snapshot(self):
        with self._lock:
            return dict(self._state), self._version

    def clear(self, pipeline: str):
        """Clear progress for a pipeline about to start."""
        with self._lock:
            if pipeline == "research":
                self._state.update(research_log="", research_kw="", research_out="", pick_log="")
            elif pipeline == "pick":
                self._state.update(pick_log="")
            elif pipeline == "keynote":
                self._state.update(keynote_log="")
            self._state["running"] = pipeline
            self._version += 1

    def finish(self):
        with self._lock:
            self._state["running"] = None
            self._version += 1


_progress = SharedProgress()

# AI provider presets: name → (base_url, description)
AI_PROVIDERS = {
    "LM Studio":       ("http://localhost:1234",  "LM Studio default"),
    "Ollama":           ("http://localhost:11434", "Ollama default"),
    "LocalAI":          ("http://localhost:8080",  "LocalAI default"),
    "vLLM":             ("http://localhost:8000",  "vLLM default"),
    "Jan":              ("http://localhost:1337",  "Jan default"),
    "llama.cpp":        ("http://localhost:8080",  "llama.cpp server"),
    "Text Gen WebUI":   ("http://localhost:5000",  "oobabooga text-generation-webui"),
    "Custom":           ("http://localhost:1234",  "Custom endpoint"),
}

# Active provider — updated at runtime from UI
_active_provider_base = "http://localhost:1234"


def get_provider_base() -> str:
    return _active_provider_base


def set_provider_base(base_url: str):
    global _active_provider_base
    _active_provider_base = base_url.rstrip("/")
PREFERRED_MODELS = [
    "qwen3.5-35b-a3b",
    "qwen3.5-27b-claude-4.6-opus-reasoning-distilled",
]

# 2-stage pipeline model routing:
#   Stage 1 (Extraction): fast bulk work — keywords, filtering, scoring
#   Stage 2 (Refinement): quality work — summaries, slides, translation
# Patterns are matched case-insensitively against available model IDs.
EXTRACTION_MODEL_PATTERN = "35b-a3b"
REFINEMENT_MODEL_PATTERN = "27b-claude"


def pick_stage_model(stage: str, available: list[str], fallback: str = "") -> str:
    """Pick the best model for a pipeline stage.

    stage: 'extraction' or 'refinement'
    available: list of loaded model IDs from LM Studio
    fallback: explicit model override from the UI (used as-is if set)

    Priority:
      1. If only one model is loaded, use it for everything.
      2. If both are loaded, route by stage.
      3. If the stage-specific model is missing, prefer the refinement model.
    """
    if fallback:
        # UI override — but still try to route if the user picked "auto"
        if fallback != "auto":
            return fallback

    if not available:
        return fallback or ""

    pattern = EXTRACTION_MODEL_PATTERN if stage == "extraction" else REFINEMENT_MODEL_PATTERN

    # Try to find the stage-specific model
    for mid in available:
        if pattern.lower() in mid.lower():
            return mid

    # Fallback: prefer refinement model (27B) > extraction model > first available
    for mid in available:
        if REFINEMENT_MODEL_PATTERN.lower() in mid.lower():
            return mid
    for mid in available:
        if EXTRACTION_MODEL_PATTERN.lower() in mid.lower():
            return mid
    return available[0]
OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "outputs")
os.makedirs(OUTPUT_DIR, exist_ok=True)
SAVED_SITES_FILE = os.path.join(os.path.dirname(__file__), ".saved_sites.json")
SITES_HISTORY_FILE = os.path.join(os.path.dirname(__file__), ".sites_history.json")
SAVED_PREFS_FILE = os.path.join(os.path.dirname(__file__), ".saved_prefs.json")

DEFAULT_PREFS = {
    "max_articles": 30,
    "slides": 10,
    "auto_keynote": False,
    "lang1": "zh-TW",
    "lang2": "English",
    "theme": "Dark",
}


def load_saved_sites() -> list[str]:
    """Load saved site URLs from disk."""
    try:
        with open(SAVED_SITES_FILE, "r") as f:
            return json.load(f)
    except (FileNotFoundError, json.JSONDecodeError):
        return []


def load_sites_history() -> list[str]:
    """Load all ever-used site URLs (history for dropdown choices)."""
    try:
        with open(SITES_HISTORY_FILE, "r") as f:
            return json.load(f)
    except (FileNotFoundError, json.JSONDecodeError):
        return []


def save_sites_history(sites: list[str]):
    """Add sites to history (never removes, only appends new ones)."""
    history = load_sites_history()
    existing = set(history)
    for s in sites:
        s = s.strip()
        if s and s not in existing:
            history.append(s)
            existing.add(s)
    with open(SITES_HISTORY_FILE, "w") as f:
        json.dump(history, f)


def load_saved_prefs() -> dict:
    """Load saved user preferences from disk."""
    try:
        with open(SAVED_PREFS_FILE, "r") as f:
            saved = json.load(f)
            return {**DEFAULT_PREFS, **saved}
    except (FileNotFoundError, json.JSONDecodeError):
        return DEFAULT_PREFS.copy()


def save_prefs(**kwargs):
    """Save user preferences to disk (merges with existing)."""
    prefs = load_saved_prefs()
    prefs.update(kwargs)
    with open(SAVED_PREFS_FILE, "w") as f:
        json.dump(prefs, f)


def save_sites(sites: list[str]):
    """Save active site URLs to disk, and add to history."""
    seen = set()
    unique = []
    for s in sites:
        s = s.strip()
        if s and s not in seen:
            seen.add(s)
            unique.append(s)
    with open(SAVED_SITES_FILE, "w") as f:
        json.dump(unique, f)
    # Also add to permanent history
    save_sites_history(unique)


def get_available_models() -> list[str]:
    """Fetch available model IDs from the active AI provider."""
    try:
        res = requests.get(f"{get_provider_base()}/v1/models", timeout=5)
        res.raise_for_status()
        return [m["id"] for m in res.json().get("data", [])]
    except Exception:
        return []


def pick_default_model(available: list[str]) -> str:
    """Return the first preferred model that is available, or the first available model."""
    for pref in PREFERRED_MODELS:
        for model_id in available:
            if pref.lower() in model_id.lower():
                return model_id
    return available[0] if available else ""


BROWSER_UA = "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/131.0.0.0 Safari/537.36"

BROWSER_HEADERS = {
    "User-Agent": BROWSER_UA,
    "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8",
    "Accept-Language": "en-US,en;q=0.9",
    "Accept-Encoding": "gzip, deflate, br",
    "Cache-Control": "no-cache",
    "Pragma": "no-cache",
    "Sec-Ch-Ua": '"Chromium";v="131", "Not_A Brand";v="24"',
    "Sec-Ch-Ua-Mobile": "?0",
    "Sec-Ch-Ua-Platform": '"macOS"',
    "Sec-Fetch-Dest": "document",
    "Sec-Fetch-Mode": "navigate",
    "Sec-Fetch-Site": "none",
    "Sec-Fetch-User": "?1",
    "Upgrade-Insecure-Requests": "1",
}


def build_session(cookies_text: str = "") -> requests.Session:
    """Build a requests.Session with optional cookies.

    Accepts either:
      - Raw cookie header string: "name1=val1; name2=val2"
      - Path to a Netscape/Mozilla cookies.txt file
    """
    session = requests.Session()
    session.headers.update(BROWSER_HEADERS)

    cookies_text = (cookies_text or "").strip()
    if not cookies_text:
        return session

    # If it looks like a file path, try loading as cookies.txt
    if os.path.isfile(cookies_text):
        jar = http.cookiejar.MozillaCookieJar(cookies_text)
        jar.load(ignore_discard=True, ignore_expires=True)
        session.cookies.update(jar)
        return session

    # Otherwise parse as "name=value; name2=value2" header string
    for pair in cookies_text.split(";"):
        pair = pair.strip()
        if "=" in pair:
            name, value = pair.split("=", 1)
            session.cookies.set(name.strip(), value.strip())

    return session


BROWSER_LOADERS = {
    "Chrome":  browser_cookie3.chrome,
    "Firefox": browser_cookie3.firefox,
    "Safari":  browser_cookie3.safari,
    "Edge":    browser_cookie3.edge,
}


def load_browser_cookies(browser_name: str, domain: str) -> str:
    """Extract cookies from an installed browser for a given domain.
    Returns a cookie header string (name=val; name2=val2).
    """
    loader = BROWSER_LOADERS.get(browser_name)
    if not loader:
        return f"Error: Unknown browser '{browser_name}'"
    try:
        jar = loader(domain_name=domain)
        pairs = [f"{c.name}={c.value}" for c in jar if domain in (c.domain or "")]
        if not pairs:
            return f"Error: No cookies found for {domain} in {browser_name}. Are you logged in?"
        return "; ".join(pairs)
    except PermissionError:
        return f"Error: Permission denied reading {browser_name} cookies. Close {browser_name} and retry, or grant Full Disk Access in System Settings > Privacy."
    except Exception as e:
        return f"Error: {e}"


def html_to_text(html: str) -> str:
    """Extract clean article text from HTML."""
    soup = BeautifulSoup(html, "lxml")
    for tag in soup.find_all(["nav", "footer", "script", "style", "aside", "iframe", "header"]):
        tag.decompose()
    article = soup.find("article") or soup.find("main") or soup.find("body")
    return article.get_text(separator="\n", strip=True) if article else soup.get_text(separator="\n", strip=True)


def extract_images_from_html(html: str, base_url: str = "", max_images: int = 10) -> list[str]:
    """Extract article images from HTML — OG image, meta images, and large article images."""
    soup = BeautifulSoup(html, "lxml")
    images = []
    seen = set()

    # 1. Open Graph image (highest quality, most relevant)
    og = soup.find("meta", property="og:image")
    if og and og.get("content"):
        images.append(og["content"])
        seen.add(og["content"])

    # 2. Twitter card image
    tw = soup.find("meta", attrs={"name": "twitter:image"})
    if tw and tw.get("content") and tw["content"] not in seen:
        images.append(tw["content"])
        seen.add(tw["content"])

    # 3. Large images from article body
    article = soup.find("article") or soup.find("main") or soup.find("body")
    if article:
        for img in article.find_all("img", src=True):
            src = img["src"]
            if src.startswith("data:"):
                continue
            src = urljoin(base_url, src) if base_url else src
            # Skip tiny icons/tracking pixels
            w = img.get("width", "")
            h = img.get("height", "")
            if w and str(w).isdigit() and int(w) < 100:
                continue
            if h and str(h).isdigit() and int(h) < 100:
                continue
            if src not in seen:
                images.append(src)
                seen.add(src)
            if len(images) >= max_images:
                break

    return images


IMG_CACHE_DIR = os.path.join(os.path.dirname(__file__), "outputs", ".img_cache")
os.makedirs(IMG_CACHE_DIR, exist_ok=True)


def download_image(url: str) -> str | None:
    """Download an image and return the local path, or None on failure."""
    try:
        res = requests.get(url, headers={"User-Agent": BROWSER_UA}, timeout=15, stream=True)
        res.raise_for_status()
        content_type = res.headers.get("Content-Type", "")
        if "image" not in content_type and not url.lower().endswith((".jpg", ".jpeg", ".png", ".webp", ".gif")):
            return None
        # Determine extension
        ext = ".jpg"
        if "png" in content_type or url.lower().endswith(".png"):
            ext = ".png"
        elif "webp" in content_type or url.lower().endswith(".webp"):
            ext = ".webp"
        import hashlib
        fname = hashlib.md5(url.encode()).hexdigest() + ext
        path = os.path.join(IMG_CACHE_DIR, fname)
        if not os.path.exists(path):
            with open(path, "wb") as f:
                for chunk in res.iter_content(8192):
                    f.write(chunk)
        return path
    except Exception:
        return None


def normalize_url(url: str) -> str:
    """Ensure URL has a scheme."""
    url = url.strip()
    if url and not url.startswith(("http://", "https://")):
        url = f"https://{url}"
    return url


def resolve_google_news_url(url: str) -> str:
    """Resolve Google News redirect URLs to actual article URLs."""
    if "news.google.com" not in url:
        return url
    try:
        from googlenewsdecoder import new_decoderv1
        result = new_decoderv1(url)
        if result and result.get("status") and result.get("decoded_url"):
            return result["decoded_url"]
    except Exception:
        pass
    return url


def _get_domain_cookies_for_playwright(url: str) -> list[dict]:
    """Extract only the target domain's cookies from Chrome — never exposes other sites."""
    parsed = urlparse(url if "://" in url else f"https://{url}")
    domain = parsed.netloc.removeprefix("www.")
    if not domain:
        return []
    try:
        jar = browser_cookie3.chrome(domain_name=domain)
        cookies = []
        for c in jar:
            # Strict match: cookie domain must be the target or a parent of the target
            # e.g. for barrons.com, allow ".barrons.com" and "www.barrons.com"
            # but never allow unrelated domains
            cookie_domain = (c.domain or "").lstrip(".")
            if cookie_domain != domain and not domain.endswith("." + cookie_domain):
                continue
            cookie = {
                "name": c.name,
                "value": c.value,
                "domain": c.domain,
                "path": c.path or "/",
            }
            if c.secure:
                cookie["secure"] = True
            if c.expires:
                cookie["expires"] = c.expires
            cookies.append(cookie)
        return cookies
    except Exception:
        return []


def _browser_fetch_html(url: str) -> str:
    """Launch a sandboxed Playwright browser with only the target domain's cookies injected."""
    from playwright.sync_api import sync_playwright

    cookies = _get_domain_cookies_for_playwright(url)
    if not cookies:
        raise RuntimeError("No cookies found for this domain")

    with sync_playwright() as p:
        browser = p.chromium.launch(
            headless=True,
            args=["--disable-blink-features=AutomationControlled"],
        )
        context = browser.new_context(
            user_agent=BROWSER_UA,
            locale="en-US",
        )
        context.add_cookies(cookies)
        page = context.new_page()
        page.goto(url, wait_until="domcontentloaded", timeout=30000)
        page.wait_for_timeout(2000)  # let JS render
        html = page.content()
        context.close()
        browser.close()

    return html


def fetch_with_browser(url: str) -> str:
    """Fetch article content using a sandboxed browser with only target-domain cookies."""
    try:
        html = _browser_fetch_html(url)
    except ImportError:
        raise RuntimeError("Playwright not installed")

    text = html_to_text(html)
    if len(text) > 200:
        return text
    raise ValueError("Browser fetch returned insufficient content")


def crawl_with_browser(url: str, max_links: int = 80) -> list[dict]:
    """Crawl links using a sandboxed browser with only target-domain cookies."""
    try:
        html = _browser_fetch_html(url)
    except Exception:
        return []

    soup = BeautifulSoup(html, "lxml")
    parsed_base = urlparse(url)
    links = []
    seen = set()

    for a in soup.find_all("a", href=True):
        href = a["href"].strip()
        full_url = urljoin(url, href)
        parsed = urlparse(full_url)
        if parsed.netloc and parsed_base.netloc.removeprefix("www.") not in parsed.netloc:
            continue
        path = parsed.path.rstrip("/")
        if not path or path == "" or len(path.split("/")) < 2:
            continue
        if full_url in seen:
            continue
        seen.add(full_url)
        title = (a.get_text(strip=True) or "")[:200]
        if title and len(title) > 10:
            links.append({"url": full_url, "title": title})

    return links[:max_links]


# Paywall indicators — if extracted text contains these, it's likely truncated/gated
_PAYWALL_MARKERS = re.compile(
    r"subscribe to continue|sign in to read|create.{0,10}account|"
    r"already a subscriber|for full access|paywall|"
    r"subscribe now|log in to read|membership required|"
    r"continue reading with|unlock this article|"
    r"this content is for subscribers",
    re.IGNORECASE,
)

# Minimum content length to consider a successful article fetch (not just a teaser)
_MIN_ARTICLE_LENGTH = 800


def _is_paywalled(text: str) -> bool:
    """Check if fetched text looks like a paywall teaser rather than full article."""
    if len(text) < _MIN_ARTICLE_LENGTH:
        return True
    if _PAYWALL_MARKERS.search(text[:2000]):
        return True
    return False


def fetch_url(url: str, session: requests.Session | None = None) -> str:
    """Fetch page content. Tries: 1) direct with session, 2) direct without, 3) browser, 4) Jina Reader."""
    url = normalize_url(url)
    url = resolve_google_news_url(url)
    best_text = ""

    # Try 1: direct fetch with session cookies
    if session and session.cookies:
        try:
            res = session.get(url, timeout=30)
            res.raise_for_status()
            text = html_to_text(res.text)
            if len(text) > 200 and not _is_paywalled(text):
                return text
            if len(text) > len(best_text):
                best_text = text
        except Exception:
            pass

    # Try 2: direct fetch without cookies (build fresh session for this domain)
    try:
        s = build_session_for_url(url)
        res = s.get(url, timeout=30)
        res.raise_for_status()
        text = html_to_text(res.text)
        if len(text) > 200 and not _is_paywalled(text):
            return text
        if len(text) > len(best_text):
            best_text = text
    except Exception:
        pass

    # Try 3: Real browser with Chrome profile (handles paywalls with active subscriptions)
    try:
        text = fetch_with_browser(url)
        if not _is_paywalled(text):
            return text
        if len(text) > len(best_text):
            best_text = text
    except Exception:
        pass

    # Try 4: Jina Reader as last resort
    try:
        res = requests.get(f"https://r.jina.ai/{url}", timeout=30)
        res.raise_for_status()
        text = res.text
        if len(text) > len(best_text):
            best_text = text
    except Exception:
        pass

    # Return best content we got, even if paywalled (partial is better than nothing)
    if best_text and len(best_text) > 200:
        return best_text

    raise ValueError(f"All fetch methods failed for {url}")


def fetch_article_images(url: str, session: requests.Session | None = None, max_images: int = 3) -> list[str]:
    """Fetch and download article images. Returns list of local file paths."""
    url = normalize_url(url)
    url = resolve_google_news_url(url)
    try:
        s = session or build_session_for_url(url)
        res = s.get(url, timeout=15)
        res.raise_for_status()
        image_urls = extract_images_from_html(res.text, url, max_images=max_images)
        local_paths = []
        for img_url in image_urls:
            path = download_image(img_url)
            if path:
                local_paths.append(path)
            if len(local_paths) >= max_images:
                break
        return local_paths
    except Exception:
        return []


def chat(prompt: str, model: str = "") -> str:
    """Send prompt to AI provider and return response text."""
    payload = {
        "model": model,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.3,
    }
    try:
        url = f"{get_provider_base()}/v1/chat/completions"
        res = requests.post(url, json=payload, timeout=300)
        res.raise_for_status()
        return res.json()["choices"][0]["message"]["content"]
    except Exception as e:
        raise ValueError(f"AI provider error: {e}")


def translate_content(content: str, language: str, model: str = "", user_prompt: str = "", lang2: str = "") -> str:
    """Translate and summarize content."""
    return chat(get_prompts().translate_and_summarize(content, language, user_prompt, lang2), model)


def _build_slide_prompt(content: str, language: str, num_slides: int, lang2: str = "", user_instruction: str = "", slide_context: str = "") -> str:
    return get_prompts().generate_slides(content, language, num_slides, lang2, user_instruction, slide_context)


def generate_slide_structure(content: str, language: str, num_slides: int, model: str = "", source_url: str = "", source_title: str = "", user_prompt: str = "", lang2: str = "") -> dict:
    """Generate structured slide JSON — batches into groups of 5 for reliability."""
    user_instruction = ""
    if user_prompt:
        user_instruction = f"""
The user's original research prompt was: "{user_prompt}"
Use this to guide the presentation angle, emphasis, and what aspects to highlight.
"""
    # Truncate content to avoid timeout
    max_content = min(10000, len(content))
    truncated = content[:max_content]

    # For small slide counts, generate in one shot
    if num_slides <= 5:
        prompt = _build_slide_prompt(truncated, language, num_slides, lang2, user_instruction)
        raw = chat(prompt, model)
        raw = raw.strip().replace("```json", "").replace("```", "").strip()
        data = json.loads(raw)
    else:
        # Batch: generate in groups of 4 slides with content chunks
        all_slides = []
        batch_size = 4
        topics_so_far = []
        first_batch_data = None
        num_batches = (num_slides + batch_size - 1) // batch_size
        # Split content across batches so each gets a portion
        chunk_size = max(2000, len(truncated) // num_batches)

        for batch_idx, batch_start in enumerate(range(0, num_slides, batch_size)):
            batch_count = min(batch_size, num_slides - batch_start)
            # Each batch gets a different chunk of content + overlap
            content_start = max(0, batch_idx * chunk_size - 500)
            content_chunk = truncated[content_start:content_start + chunk_size + 500]

            context = ""
            if topics_so_far:
                context = f"Previous slides covered: {', '.join(topics_so_far)}. Generate {batch_count} NEW slides on DIFFERENT angles."

            prompt = _build_slide_prompt(content_chunk, language, batch_count, lang2, user_instruction, context)
            try:
                raw = chat(prompt, model)
                raw = raw.strip().replace("```json", "").replace("```", "").strip()
                batch_data = json.loads(raw)
            except Exception:
                continue
            batch_slides = batch_data.get("slides", [])
            all_slides.extend(batch_slides)
            if first_batch_data is None:
                first_batch_data = batch_data

            for s in batch_slides:
                heading = s.get("heading_primary", s.get("heading_en", s.get("heading", "")))
                if heading:
                    topics_so_far.append(heading[:30])

        data = first_batch_data if first_batch_data else {"slides": all_slides}
        data["slides"] = all_slides

    # Attach source info
    data = data if isinstance(data, dict) else {"slides": all_slides}
    if source_url:
        data["source_url"] = source_url
        data["source_title"] = source_title or source_url
    return data


def add_hyperlink(paragraph, url, text, font_size, color, slide_part=None):
    """Add a clickable hyperlink run to a paragraph."""
    from pptx.oxml.ns import qn
    run = paragraph.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.underline = True

    if not slide_part:
        # Walk up to find the slide part
        obj = paragraph._parent
        while obj is not None:
            if hasattr(obj, 'part'):
                slide_part = obj.part
                break
            obj = getattr(obj, '_parent', None) or getattr(obj, 'parent', None)

    if not slide_part:
        return  # Can't create hyperlink without slide part

    # Register the relationship
    rel_type = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink"
    rel = slide_part.relate_to(url, rel_type, is_external=True)
    # rel can be a string rId or an object with .rId
    rId = rel if isinstance(rel, str) else rel.rId

    # Create the hyperlink XML on the run
    r_elem = run._r
    hlinkClick = r_elem.makeelement(qn("a:hlinkClick"), {})
    hlinkClick.set(qn("r:id"), rId)
    r_elem.append(hlinkClick)


def build_pptx(data: dict, theme: str, issues: list[str] | None = None, cover_image: str | None = None) -> str:
    """Build .pptx file from slide data — dual-language with source links."""
    themes = {
        "Dark":  {"bg": RGBColor(0x1A, 0x1A, 0x2E), "title": RGBColor(0xE9, 0x4F, 0x37), "text": RGBColor(0xFF, 0xFF, 0xFF), "sub": RGBColor(0xAA, 0xAA, 0xAA), "link": RGBColor(0x4F, 0xC3, 0xF7)},
        "Light": {"bg": RGBColor(0xFA, 0xFA, 0xFA), "title": RGBColor(0x22, 0x22, 0x22),  "text": RGBColor(0x33, 0x33, 0x33), "sub": RGBColor(0x77, 0x77, 0x77), "link": RGBColor(0x1A, 0x73, 0xE8)},
        "Blue":  {"bg": RGBColor(0x0F, 0x3D, 0x66), "title": RGBColor(0x4F, 0xC3, 0xF7), "text": RGBColor(0xFF, 0xFF, 0xFF), "sub": RGBColor(0xAA, 0xCC, 0xDD), "link": RGBColor(0x8B, 0xD9, 0xFF)},
    }
    colors = themes.get(theme, themes["Dark"])

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    source_url = data.get("source_url", "")
    all_sources = data.get("all_sources", [source_url] if source_url else [])

    def set_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    has_secondary = bool(data.get("title_secondary"))

    # --- Title Slide ---
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(title_slide, colors["bg"])

    # Primary title
    txBox = title_slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title_primary", data.get("title_zh", data.get("title", "")))
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = colors["title"]

    # Secondary title
    if has_secondary:
        enBox = title_slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(10), Inches(1))
        entf = enBox.text_frame
        enp = entf.paragraphs[0]
        enp.text = data.get("title_secondary", data.get("title_en", ""))
        enp.font.size = Pt(28)
        enp.font.color.rgb = colors["sub"]

    # Subtitle
    sub_y = 4.5 if has_secondary else 3.2
    subBox = title_slide.shapes.add_textbox(Inches(1.5), Inches(sub_y), Inches(10), Inches(1))
    stf = subBox.text_frame
    sp = stf.paragraphs[0]
    sp.text = data.get("subtitle_primary", data.get("subtitle_zh", data.get("subtitle", "")))
    sp.font.size = Pt(20)
    sp.font.color.rgb = colors["sub"]
    sec_sub = data.get("subtitle_secondary", data.get("subtitle_en", ""))
    if sec_sub:
        sp2 = stf.add_paragraph()
        sp2.text = sec_sub
        sp2.font.size = Pt(18)
        sp2.font.color.rgb = colors["sub"]

    # Cover image on title slide (lower right, below text)
    if cover_image and os.path.exists(cover_image):
        try:
            from PIL import Image
            with Image.open(cover_image) as img:
                img_w, img_h = img.size
            aspect = img_w / img_h if img_h else 1
            max_w, max_h = 3.5, 2.5
            if aspect > max_w / max_h:
                w = min(max_w, img_w / 96)
                h = w / aspect
            else:
                h = min(max_h, img_h / 96)
                w = h * aspect
            # Position: bottom-right corner with padding
            left = Inches(13.33 - w - 0.6)
            top = Inches(7.5 - h - 0.6)
            title_slide.shapes.add_picture(cover_image, left, top, Inches(w), Inches(h))
        except ImportError:
            title_slide.shapes.add_picture(cover_image, Inches(9.0), Inches(4.5), Inches(3.5), Inches(2.2))
        except Exception:
            pass

    # Source links on title slide
    if all_sources:
        src_y = 5.2 if len(all_sources) > 1 else 5.8
        srcBox = title_slide.shapes.add_textbox(Inches(1.5), Inches(src_y), Inches(8), Inches(1.8))
        srctf = srcBox.text_frame
        srctf.word_wrap = True
        for si, src_url in enumerate(all_sources):
            src_domain = urlparse(normalize_url(src_url)).netloc
            sp = srctf.paragraphs[0] if si == 0 else srctf.add_paragraph()
            add_hyperlink(sp, src_url, f"📎 {src_domain}", Pt(11), colors["link"], slide_part=title_slide.part)
            sp.space_after = Pt(2)

    # --- Summary Slide ---
    summary_text = data.get("summary", "")
    if summary_text:
        # Split into lines first, then paginate by line count
        all_lines = [l.strip() for l in summary_text.split("\n") if l.strip()]
        max_lines_per_slide = 14
        summary_pages = []
        for i in range(0, len(all_lines), max_lines_per_slide):
            summary_pages.append(all_lines[i:i + max_lines_per_slide])

        # If only one page but too many chars, split by char count
        if len(summary_pages) == 1 and len(summary_text) > 700:
            summary_pages = []
            remaining = summary_text
            while remaining:
                if len(remaining) <= 700:
                    summary_pages.append(remaining.split("\n"))
                    break
                cut = remaining.rfind("。", 0, 700)
                if cut < 0:
                    cut = remaining.rfind(". ", 0, 700)
                if cut < 0:
                    cut = remaining.rfind("\n", 0, 700)
                if cut < 0:
                    cut = 700
                summary_pages.append([l.strip() for l in remaining[:cut + 1].split("\n") if l.strip()])
                remaining = remaining[cut + 1:].strip()

        for pg_idx, pg_lines in enumerate(summary_pages):
            sum_slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_bg(sum_slide, colors["bg"])

            stitle = "Executive Summary"
            if len(summary_pages) > 1:
                stitle += f" ({pg_idx + 1}/{len(summary_pages)})"
            shBox = sum_slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.7))
            shtf = shBox.text_frame
            shp = shtf.paragraphs[0]
            shp.text = stitle
            shp.font.size = Pt(28)
            shp.font.bold = True
            shp.font.color.rgb = colors["title"]

            sbBox = sum_slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.5), Inches(6.0))
            sbtf = sbBox.text_frame
            sbtf.word_wrap = True
            first = True
            for line in pg_lines:
                if not line:
                    continue
                sp = sbtf.paragraphs[0] if first else sbtf.add_paragraph()
                first = False
                sp.text = line
                sp.font.size = Pt(14)
                sp.font.color.rgb = colors["text"]
                sp.space_after = Pt(3)

    # --- Content Slides ---
    for slide_idx, slide_data in enumerate(data["slides"]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, colors["bg"])

        # Primary heading
        hBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.7))
        htf = hBox.text_frame
        hp = htf.paragraphs[0]
        hp.text = slide_data.get("heading_primary", slide_data.get("heading_zh", slide_data.get("heading", "")))
        hp.font.size = Pt(32)
        hp.font.bold = True
        hp.font.color.rgb = colors["title"]

        # Secondary heading
        if has_secondary:
            heBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.5))
            hetf = heBox.text_frame
            hep = hetf.paragraphs[0]
            hep.text = slide_data.get("heading_secondary", slide_data.get("heading_en", ""))
            hep.font.size = Pt(20)
            hep.font.color.rgb = colors["sub"]

        # Bullets
        pri_bullets = slide_data.get("bullets_primary", slide_data.get("bullets_zh", slide_data.get("bullets", [])))
        sec_bullets = slide_data.get("bullets_secondary", slide_data.get("bullets_en", []))

        # Primary bullets (left column if dual, full width if single)
        bullet_top = Inches(1.7) if has_secondary else Inches(1.5)
        bullet_w = Inches(5.5) if has_secondary else Inches(11)
        bBox = slide.shapes.add_textbox(Inches(0.8), bullet_top, bullet_w, Inches(4.5))
        btf = bBox.text_frame
        btf.word_wrap = True
        for i, bullet in enumerate(pri_bullets):
            bp = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            bp.text = f"• {bullet}"
            bp.font.size = Pt(18)
            bp.font.color.rgb = colors["text"]
            bp.space_after = Pt(6)

        # Secondary bullets (right column, only if dual-language)
        if has_secondary and sec_bullets:
            eBox = slide.shapes.add_textbox(Inches(6.8), bullet_top, Inches(5.5), Inches(4.5))
            etf = eBox.text_frame
            etf.word_wrap = True
            for i, bullet in enumerate(sec_bullets):
                ep = etf.paragraphs[0] if i == 0 else etf.add_paragraph()
                ep.text = f"• {bullet}"
                ep.font.size = Pt(16)
                ep.font.color.rgb = colors["sub"]
                ep.space_after = Pt(6)

        # Source link at bottom — rotate through all sources
        if all_sources:
            src = all_sources[slide_idx % len(all_sources)]
            src_domain = urlparse(normalize_url(src)).netloc
            linkBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11), Inches(0.4))
            ltf = linkBox.text_frame
            lp = ltf.paragraphs[0]
            add_hyperlink(lp, src, f"📎 {src_domain}", Pt(10), colors["link"], slide_part=slide.part)

    # --- Issues/Warnings Slides (paginated) ---
    if issues:
        issues_per_page = 12
        for page_start in range(0, len(issues), issues_per_page):
            page_issues = issues[page_start:page_start + issues_per_page]
            page_num = page_start // issues_per_page + 1
            total_pages = (len(issues) + issues_per_page - 1) // issues_per_page

            issue_slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_bg(issue_slide, colors["bg"])

            title_text = "⚠️ Issues & Warnings"
            if total_pages > 1:
                title_text += f" ({page_num}/{total_pages})"

            ihBox = issue_slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.7))
            ihtf = ihBox.text_frame
            ihp = ihtf.paragraphs[0]
            ihp.text = title_text
            ihp.font.size = Pt(28)
            ihp.font.bold = True
            ihp.font.color.rgb = RGBColor(0xE9, 0x4F, 0x37)

            ibBox = issue_slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5))
            ibtf = ibBox.text_frame
            ibtf.word_wrap = True
            for i, issue in enumerate(page_issues):
                ip = ibtf.paragraphs[0] if i == 0 else ibtf.add_paragraph()
                ip.text = f"• {issue}"
                ip.font.size = Pt(13)
                ip.font.color.rgb = RGBColor(0xFF, 0xAA, 0x00)
                ip.space_after = Pt(3)

    # Build descriptive filename: date + English title + translated title
    timestamp = datetime.now().strftime("%Y%m%d")
    title_pri = data.get("title_primary", data.get("title_zh", data.get("title", "")))
    title_sec = data.get("title_secondary", data.get("title_en", ""))
    # Figure out which is English — secondary is usually English when primary is translated
    # If primary looks like English (ASCII), swap
    is_pri_english = all(ord(c) < 128 for c in title_pri.replace(" ", "")) if title_pri else False
    if is_pri_english:
        title_en = title_pri
        title_translated = title_sec
    else:
        title_en = title_sec
        title_translated = title_pri

    def _sanitize(s: str, max_len: int = 30) -> str:
        """Clean string for use in filename."""
        s = re.sub(r'[\\/:*?"<>|&\'()!@#$%^{};\[\]]', '', s)
        s = re.sub(r'\s+', '_', s.strip())
        return s[:max_len].rstrip('_')

    parts = [timestamp]
    if title_en:
        parts.append(_sanitize(title_en))
    if title_translated and title_translated != title_en:
        parts.append(_sanitize(title_translated))
    filename = "_".join(parts) or f"keynote_{timestamp}"

    path = os.path.join(OUTPUT_DIR, f"{filename}.pptx")
    prs.save(path)
    return path


def extract_keywords(prompt: str, model: str = "") -> list[str]:
    """Use AI to break down a prompt into broad search keywords."""
    raw = chat(get_prompts().extract_keywords(prompt), model)
    raw = raw.strip().replace("```json", "").replace("```", "").strip()
    return json.loads(raw)


RSS_FEEDS = {
    # --- Finance / News ---
    "marketwatch.com": [
        "https://www.marketwatch.com/rss/topstories",
        "https://www.marketwatch.com/rss/marketpulse",
    ],
    "wsj.com": [
        "https://feeds.a.dj.com/rss/RSSWorldNews.xml",
        "https://feeds.a.dj.com/rss/WSJcomUSBusiness.xml",
        "https://feeds.a.dj.com/rss/RSSMarketsMain.xml",
    ],
    "barrons.com": [
        "https://feeds.barrons.com/barrons/review",
        "https://feeds.barrons.com/barrons/market_lab",
    ],
    "economist.com": [
        "https://www.economist.com/finance-and-economics/rss.xml",
        "https://www.economist.com/business/rss.xml",
        "https://www.economist.com/international/rss.xml",
        "https://www.economist.com/leaders/rss.xml",
    ],
    # --- Tech ---
    "techcrunch.com": [
        "https://techcrunch.com/feed/",
    ],
    "arstechnica.com": [
        "https://feeds.arstechnica.com/arstechnica/index",
    ],
    "theverge.com": [
        "https://www.theverge.com/rss/index.xml",
    ],
    "wired.com": [
        "https://www.wired.com/feed/rss",
    ],
    "thenextweb.com": [
        "https://thenextweb.com/feed",
    ],
    "venturebeat.com": [
        "https://venturebeat.com/feed/",
    ],
    "9to5mac.com": [
        "https://9to5mac.com/feed/",
    ],
    "engadget.com": [
        "https://www.engadget.com/rss.xml",
    ],
    "zdnet.com": [
        "https://www.zdnet.com/news/rss.xml",
    ],
    "tomshardware.com": [
        "https://www.tomshardware.com/feeds/all",
    ],
}

# Sites that need Google News RSS as proxy (strong bot detection, no RSS)
GOOGLE_NEWS_FALLBACK = {"reuters.com", "investopedia.com"}

SAVED_CATEGORIES_FILE = os.path.join(os.path.dirname(__file__), ".saved_categories.json")

DEFAULT_CATEGORIES = {
    "Finance": [
        "https://www.marketwatch.com",
        "https://www.wsj.com",
        "https://www.barrons.com",
        "https://www.economist.com",
        "https://www.reuters.com",
        "https://www.investopedia.com",
    ],
    "Tech": [
        "https://techcrunch.com",
        "https://arstechnica.com",
        "https://www.theverge.com",
        "https://www.wired.com",
        "https://thenextweb.com",
        "https://venturebeat.com",
        "https://9to5mac.com",
        "https://www.engadget.com",
        "https://www.zdnet.com",
        "https://www.tomshardware.com",
    ],
}


def load_categories() -> dict:
    """Load saved categories from disk, merged with defaults."""
    try:
        with open(SAVED_CATEGORIES_FILE, "r") as f:
            saved = json.load(f)
            # Merge defaults for any missing categories
            merged = {**DEFAULT_CATEGORIES}
            merged.update(saved)
            return merged
    except (FileNotFoundError, json.JSONDecodeError):
        return DEFAULT_CATEGORIES.copy()


def save_categories(categories: dict):
    """Save categories to disk."""
    with open(SAVED_CATEGORIES_FILE, "w") as f:
        json.dump(categories, f, indent=2)


def get_all_suggested_sites(categories: dict) -> list[str]:
    """Flatten all category sites into a single list for dropdown choices."""
    all_sites = []
    seen = set()
    for sites in categories.values():
        for s in sites:
            if s not in seen:
                all_sites.append(s)
                seen.add(s)
    return all_sites


def crawl_rss_feeds(feed_urls: list[str], max_links: int = 80) -> list[dict]:
    """Fetch article links from RSS feeds."""
    links = []
    seen = set()
    for feed_url in feed_urls:
        try:
            res = requests.get(feed_url, headers={"User-Agent": BROWSER_UA}, timeout=15)
            res.raise_for_status()
            soup = BeautifulSoup(res.content, "lxml-xml")
            for item in soup.find_all("item"):
                title_tag = item.find("title")
                link_tag = item.find("link")
                if not title_tag or not link_tag:
                    continue
                title = title_tag.get_text(strip=True)
                url = link_tag.get_text(strip=True)
                if url in seen or not title or len(title) < 10:
                    continue
                seen.add(url)
                links.append({"url": url, "title": title})
                if len(links) >= max_links:
                    return links
        except Exception:
            continue
    return links


def crawl_google_news(domain: str, max_links: int = 80) -> list[dict]:
    """Use Google News RSS as proxy for sites that block direct crawl."""
    feed_url = f"https://news.google.com/rss/search?q=site:{domain}&hl=en-US&gl=US&ceid=US:en"
    try:
        res = requests.get(feed_url, headers={"User-Agent": BROWSER_UA}, timeout=15)
        res.raise_for_status()
    except Exception:
        return []

    soup = BeautifulSoup(res.content, "lxml-xml")
    links = []
    seen = set()
    for item in soup.find_all("item"):
        title_tag = item.find("title")
        link_tag = item.find("link")
        if not title_tag or not link_tag:
            continue
        title = title_tag.get_text(strip=True)
        # Clean title — Google appends " - Source Name"
        title = re.sub(r'\s*-\s*[^-]+$', '', title).strip()
        gn_url = link_tag.get_text(strip=True)
        # Resolve Google redirect to actual article URL
        url = resolve_google_news_url(gn_url)
        if url in seen or not title or len(title) < 10:
            continue
        seen.add(url)
        links.append({"url": url, "title": title})
        if len(links) >= max_links:
            break
    return links


# Domains where article content can't be fetched (paywall + cookie encryption)
# Titles from these are still used for AI filtering, but deep-scan will skip them

# Search engines — when user adds these as a "site", we search instead of crawl
# Note: Google blocks scraping; DuckDuckGo and Bing work reliably
SEARCH_ENGINES = {
    "google.com": "https://www.google.com/search?q={query}&num={num}",
    "bing.com": "https://www.bing.com/search?q={query}&count={num}",
    "duckduckgo.com": "https://html.duckduckgo.com/html/?q={query}",
}

def search_engine_crawl(engine_domain: str, query: str, max_links: int = 80) -> list[dict]:
    """Use a search engine to find articles related to a query."""
    links = []
    seen = set()

    # Build search URL
    template = None
    for domain, url_template in SEARCH_ENGINES.items():
        if domain in engine_domain:
            template = url_template
            break
    if not template:
        return []

    search_url = template.format(query=requests.utils.quote(query), num=min(max_links, 30))

    try:
        headers = {
            "User-Agent": BROWSER_UA,
            "Accept": "text/html,application/xhtml+xml",
            "Accept-Language": "en-US,en;q=0.9",
        }
        res = requests.get(search_url, headers=headers, timeout=15)
        res.raise_for_status()
    except Exception:
        return []

    soup = BeautifulSoup(res.text, "lxml")

    def _extract_real_url(href: str) -> str:
        """Extract real URL from search engine redirect wrappers."""
        # DuckDuckGo wraps: //duckduckgo.com/l/?uddg=https%3A%2F%2F...
        if "uddg=" in href:
            from urllib.parse import unquote, parse_qs
            parsed_q = parse_qs(urlparse(href).query)
            if "uddg" in parsed_q:
                return unquote(parsed_q["uddg"][0])
        # Google wraps: /url?q=https://...
        if "/url?" in href and "q=" in href:
            from urllib.parse import unquote, parse_qs
            parsed_q = parse_qs(urlparse(href).query)
            if "q" in parsed_q:
                return unquote(parsed_q["q"][0])
        return href

    # Google results
    if "google.com" in engine_domain:
        for div in soup.find_all("div", class_="g"):
            a = div.find("a", href=True)
            h3 = div.find("h3")
            if not a or not h3:
                continue
            url = _extract_real_url(a["href"])
            title = h3.get_text(strip=True)
            if not url.startswith("http") or "google.com" in url:
                continue
            if url not in seen and title and len(title) > 5:
                seen.add(url)
                links.append({"url": url, "title": title})
        # Fallback: if no div.g results, try broader search
        if not links:
            for a in soup.find_all("a", href=True):
                url = _extract_real_url(a["href"])
                h3 = a.find("h3")
                if not h3:
                    continue
                title = h3.get_text(strip=True)
                if not url.startswith("http") or "google.com" in url:
                    continue
                if url not in seen and title and len(title) > 5:
                    seen.add(url)
                    links.append({"url": url, "title": title})

    # Bing results
    elif "bing.com" in engine_domain:
        for li in soup.find_all("li", class_="b_algo"):
            a = li.find("a", href=True)
            if not a:
                continue
            url = _extract_real_url(a["href"])
            title = a.get_text(strip=True)
            if not url.startswith("http") or "bing.com" in url:
                continue
            if url not in seen and title and len(title) > 5:
                seen.add(url)
                links.append({"url": url, "title": title})

    # DuckDuckGo results
    elif "duckduckgo.com" in engine_domain:
        for a in soup.find_all("a", class_="result__a"):
            raw_href = a.get("href", "")
            url = _extract_real_url(raw_href)
            title = a.get_text(strip=True)
            if not url.startswith("http") or "duckduckgo.com" in url:
                continue
            if url not in seen and title and len(title) > 5:
                seen.add(url)
                links.append({"url": url, "title": title})

    return links[:max_links]


def is_search_engine(domain: str) -> bool:
    """Check if a domain is a search engine."""
    return any(se in domain for se in SEARCH_ENGINES)


def crawl_site_links(base_url: str, session: requests.Session | None = None, max_links: int = 80) -> list[dict]:
    """Crawl a website. Falls back to RSS/Google News if direct crawl fails."""
    base_url = normalize_url(base_url)
    parsed = urlparse(base_url)
    domain = parsed.netloc.removeprefix("www.")

    try:
        s = session or requests.Session()
        for k, v in BROWSER_HEADERS.items():
            s.headers.setdefault(k, v)
        # If session has no cookies, try auto-loading for this domain
        if not s.cookies:
            fresh = build_session_for_url(base_url)
            if fresh.cookies:
                s.cookies.update(fresh.cookies)
        res = s.get(base_url, timeout=15)
        res.raise_for_status()
    except Exception as e:
        # Fallback 1: Real browser with Chrome profile (paywall sites with subscriptions)
        browser_links = crawl_with_browser(base_url, max_links)
        if browser_links:
            return browser_links
        # Fallback 2: RSS feeds
        for rss_domain, feeds in RSS_FEEDS.items():
            if rss_domain in domain:
                rss_links = crawl_rss_feeds(feeds, max_links)
                if rss_links:
                    return rss_links
        # Fallback 3: Google News RSS proxy
        for gn_domain in GOOGLE_NEWS_FALLBACK:
            if gn_domain in domain:
                gn_links = crawl_google_news(gn_domain, max_links)
                if gn_links:
                    return gn_links
        raise ValueError(f"Failed to fetch {base_url}: {e}")

    soup = BeautifulSoup(res.text, "lxml")
    parsed_base = urlparse(base_url)
    links = []
    seen = set()

    for a in soup.find_all("a", href=True):
        href = a["href"].strip()
        full_url = urljoin(base_url, href)
        parsed = urlparse(full_url)

        if parsed.netloc != parsed_base.netloc:
            continue
        # Deduplicate by path (ignore query params like ?mod=...)
        canon = f"{parsed.scheme}://{parsed.netloc}{parsed.path}"
        if canon in seen:
            continue
        path_parts = [pt for pt in parsed.path.split("/") if pt]
        if len(path_parts) < 2:
            continue
        skip_paths = {"login", "signup", "subscribe", "account", "search", "video", "podcasts", "#"}
        if any(sk in parsed.path.lower() for sk in skip_paths):
            continue

        title = a.get_text(strip=True)
        if not title or len(title) < 10:
            continue

        seen.add(canon)
        links.append({"url": canon, "title": title})
        if len(links) >= max_links:
            break

    return links


# Patterns that indicate a section/category/topic page (not an individual article)
_SECTION_PATTERNS = re.compile(
    r"/(topic|topics|category|categories|section|sections|tag|tags|"
    r"latest|news|opinion|analysis|markets|technology|tech|science|"
    r"business|finance|economy|world|politics|industry|ai|crypto|"
    r"startups|reviews|features|columns|archive)(/|$)", re.IGNORECASE
)

# Patterns that indicate an individual article (should NOT be followed for depth 2)
_ARTICLE_PATTERNS = re.compile(
    r"/\d{4}/\d{2}/|"           # date-based paths: /2024/03/
    r"/[a-z0-9-]{30,}$|"        # long slugs (article titles)
    r"/\d{5,}$|"                # numeric article IDs
    r"\.(pdf|jpg|png|mp4)$",    # file extensions
    re.IGNORECASE
)


def _is_section_page(url: str, base_domain: str) -> bool:
    """Detect if a URL is likely a section/category page worth crawling deeper."""
    parsed = urlparse(url)
    # Must be same domain
    if base_domain not in parsed.netloc:
        return False
    path = parsed.path.rstrip("/")
    # Section pages typically have short paths (1-2 segments)
    parts = [p for p in path.split("/") if p]
    if len(parts) > 3 or len(parts) < 1:
        return False
    # Match section patterns
    if _SECTION_PATTERNS.search(path):
        return True
    # Short path with 1-2 segments that don't look like articles
    if len(parts) <= 2 and not _ARTICLE_PATTERNS.search(path):
        return True
    return False


def _extract_links_from_html(html: str, base_url: str, seen: set, max_links: int = 80) -> list[dict]:
    """Extract article links from HTML, deduplicating against seen set."""
    soup = BeautifulSoup(html, "lxml")
    parsed_base = urlparse(base_url)
    links = []

    for a in soup.find_all("a", href=True):
        href = a["href"].strip()
        full_url = urljoin(base_url, href)
        parsed = urlparse(full_url)

        if parsed.netloc != parsed_base.netloc:
            continue
        canon = f"{parsed.scheme}://{parsed.netloc}{parsed.path}"
        if canon in seen:
            continue
        path_parts = [pt for pt in parsed.path.split("/") if pt]
        if len(path_parts) < 2:
            continue
        skip_paths = {"login", "signup", "subscribe", "account", "search", "video", "podcasts", "#"}
        if any(sk in parsed.path.lower() for sk in skip_paths):
            continue
        title = a.get_text(strip=True)
        if not title or len(title) < 10:
            continue

        seen.add(canon)
        links.append({"url": canon, "title": title})
        if len(links) >= max_links:
            break

    return links


def crawl_site_deep(base_url: str, keywords: list[str], session: requests.Session | None = None,
                    max_links: int = 150, max_sub_pages: int = 10, log_fn=None) -> list[dict]:
    """Depth-2 crawl: crawl main page, then selectively follow section/category pages.

    Only follows pages that match section patterns or contain keywords in their URL/text.
    Caps sub-page crawls to max_sub_pages to keep resource usage bounded.
    """
    base_url = normalize_url(base_url)
    parsed = urlparse(base_url)
    domain = parsed.netloc.removeprefix("www.")

    s = session or requests.Session()
    for k, v in BROWSER_HEADERS.items():
        s.headers.setdefault(k, v)
    if not s.cookies:
        fresh = build_session_for_url(base_url)
        if fresh.cookies:
            s.cookies.update(fresh.cookies)

    # Depth 1: crawl main page
    try:
        res = s.get(base_url, timeout=15)
        res.raise_for_status()
        main_html = res.text
    except Exception:
        # Try browser fallback
        try:
            main_html = _browser_fetch_html(base_url)
        except Exception:
            # Fall back to regular crawl (RSS etc)
            return crawl_site_links(base_url, session=session, max_links=max_links)

    seen = set()
    all_links = _extract_links_from_html(main_html, base_url, seen, max_links)

    if log_fn:
        log_fn(f"    📄 Depth 1: {len(all_links)} links")

    # Identify section pages to follow
    soup = BeautifulSoup(main_html, "lxml")
    section_urls = []
    kw_lower = [k.lower() for k in keywords] if keywords else []

    for a in soup.find_all("a", href=True):
        href = a["href"].strip()
        full_url = urljoin(base_url, href)
        canon = f"{urlparse(full_url).scheme}://{urlparse(full_url).netloc}{urlparse(full_url).path}"

        if canon == base_url.rstrip("/") or canon in seen:
            continue

        if _is_section_page(full_url, domain):
            # Prioritize sections that match keywords
            link_text = (a.get_text(strip=True) + " " + urlparse(full_url).path).lower()
            keyword_match = any(kw in link_text for kw in kw_lower) if kw_lower else False
            section_urls.append((full_url, keyword_match))

    # Sort: keyword-matching sections first, then others
    section_urls.sort(key=lambda x: (not x[1], x[0]))
    section_urls = [url for url, _ in section_urls[:max_sub_pages]]

    if log_fn and section_urls:
        log_fn(f"    🔍 Depth 2: following {len(section_urls)} section pages...")

    # Depth 2: crawl selected section pages
    for sub_url in section_urls:
        if len(all_links) >= max_links:
            break
        try:
            res = s.get(sub_url, timeout=15)
            res.raise_for_status()
            sub_html = res.text
        except Exception:
            try:
                sub_html = _browser_fetch_html(sub_url)
            except Exception:
                continue

        remaining = max_links - len(all_links)
        new_links = _extract_links_from_html(sub_html, sub_url, seen, remaining)
        all_links.extend(new_links)

        if log_fn and new_links:
            section_name = urlparse(sub_url).path.rstrip("/").split("/")[-1] or "index"
            log_fn(f"    📄 /{section_name}: +{len(new_links)} links")

    return all_links


def score_text(text: str, keywords: list[str]) -> float:
    """Score text against keywords (0-1). Matches both full phrases and individual words."""
    if not keywords:
        return 0
    text_lower = text.lower()
    score = 0
    for kw in keywords:
        kw_lower = kw.lower()
        if kw_lower in text_lower:
            # Full keyword match — full point
            score += 1
        else:
            # Partial: check individual words within the keyword
            words = kw_lower.split()
            if len(words) > 1:
                word_hits = sum(1 for w in words if w in text_lower)
                score += word_hits / len(words)
    return score / len(keywords)


def score_article(title: str, keywords: list[str]) -> float:
    """Score an article title against keywords."""
    return score_text(title, keywords)


def ai_filter_titles(titles: list[dict], prompt: str, model: str = "", keywords: list[str] = None, batch_size: int = 40) -> list[dict]:
    """Pre-filter by keywords, then use AI to judge relevance in batches."""
    # Step 1: fast keyword pre-filter — always run, keep top 80
    max_to_ai = 80
    if keywords:
        scored = []
        for t in titles:
            s = score_text(t["title"], keywords)
            if s > 0.05:  # Require meaningful match, not just 1 common word
                scored.append((s, t))
        scored.sort(key=lambda x: x[0], reverse=True)
        titles = [t for _, t in scored[:max_to_ai]]
    elif len(titles) > max_to_ai:
        titles = titles[:max_to_ai]

    # Step 2: batch AI filtering (smaller batches for reliability)
    all_results = []
    for batch_start in range(0, len(titles), batch_size):
        batch = titles[batch_start:batch_start + batch_size]
        try:
            batch_results = _ai_filter_batch(batch, prompt, model)
            all_results.extend(batch_results)
        except Exception:
            # If a batch fails, skip it rather than crashing everything
            continue

    return all_results


def _ai_filter_batch(titles: list[dict], prompt: str, model: str = "") -> list[dict]:
    """AI judges a single batch of titles."""
    title_list = "\n".join(f"{i}. {t['title']}" for i, t in enumerate(titles))
    raw = chat(get_prompts().filter_titles(prompt, title_list), model)
    raw = raw.strip().replace("```json", "").replace("```", "").strip()
    try:
        judgments = json.loads(raw)
    except json.JSONDecodeError:
        return []

    results = []
    for j in judgments:
        idx = j.get("index", -1)
        if 0 <= idx < len(titles):
            results.append({
                **titles[idx],
                "score": j.get("relevance", 0.5),
                "reason": j.get("reason", ""),
            })
    return results


def ai_score_article(url: str, prompt: str, session: requests.Session | None = None, model: str = "") -> dict:
    """Fetch article content and use AI to judge relevance to user's topic."""
    try:
        content = fetch_url(url, session)[:4000]
    except Exception:
        return {"score": 0, "reason": "Failed to fetch"}

    try:
        raw = chat(get_prompts().score_article(prompt, content), model)
        raw = raw.strip().replace("```json", "").replace("```", "").strip()
        result = json.loads(raw)
        return {"score": result.get("relevance", 0), "reason": result.get("reason", "")}
    except Exception:
        return {"score": 0, "reason": "AI scoring failed"}


def build_session_for_url(url: str, cookies_text: str = "") -> requests.Session:
    """Build a session with cookies specific to the URL's domain."""
    session = requests.Session()
    session.headers.update(BROWSER_HEADERS)
    # Set Referer to the site's homepage so it looks like in-site navigation
    parsed = urlparse(url if "://" in url else f"https://{url}")
    session.headers["Referer"] = f"{parsed.scheme}://{parsed.netloc}/"
    # Increase header limit for sites like Yahoo that send excessive Set-Cookie headers
    import http.client
    http.client._MAXHEADERS = 1000

    # If manual cookies provided, use those
    if cookies_text and cookies_text.strip():
        for pair in cookies_text.split(";"):
            pair = pair.strip()
            if "=" in pair:
                name, value = pair.split("=", 1)
                session.cookies.set(name.strip(), value.strip())
        return session

    # Auto-load cookies for this specific domain
    parsed = urlparse(url if "://" in url else f"https://{url}")
    domain = parsed.netloc.removeprefix("www.")
    if domain:
        try:
            result = load_browser_cookies("Chrome", domain)
            if not result.startswith("Error:"):
                for pair in result.split(";"):
                    pair = pair.strip()
                    if "=" in pair:
                        name, value = pair.split("=", 1)
                        session.cookies.set(name.strip(), value.strip())
        except Exception:
            pass  # Cookie loading failed silently, proceed without
    return session


_research_issues = []  # Shared between research and pipeline
_research_results = []  # Store latest results for topic grouping
_active_stop_flag = None  # Global stop flag for the current running pipeline

def run_research(prompt, site_urls, max_articles, model, cookies_text, crawl_deep=False, stop_flag=None):
    """Research pipeline: keywords → crawl multiple sites → AI filters titles → deep scan."""
    global _research_issues
    _research_issues = []
    logs = []
    no_update = gr.update()
    if stop_flag is None:
        stop_flag = StopFlag()  # no-op flag if none provided

    # Normalize site_urls to a list
    if isinstance(site_urls, str):
        url_list = [u.strip() for u in site_urls.split(";") if u.strip()]
    elif isinstance(site_urls, list):
        url_list = [u.strip() for u in site_urls if u and u.strip()]
    else:
        url_list = []

    if not url_list:
        yield "❌ No site URLs provided", "", "", no_update
        return

    # Save sites for future use
    save_sites(url_list)

    # Resolve stage-specific models
    available = get_available_models()
    ext_model = pick_stage_model("extraction", available, model)
    ref_model = pick_stage_model("refinement", available, model)
    if ext_model != ref_model:
        logs.append(f"🔀 2-stage pipeline: extraction → {ext_model.split('/')[-1]}, refinement → {ref_model.split('/')[-1]}")
    else:
        logs.append(f"🤖 Model: {ext_model.split('/')[-1]}")

    try:
        stop_flag.check()
        logs.append("🔍 Analyzing prompt with AI...")
        yield "\n".join(logs), "", "", no_update

        keywords = extract_keywords(prompt, ext_model)
        kw_display = ", ".join(keywords)
        logs.append(f"🔑 Keywords: {kw_display}")
        yield "\n".join(logs), kw_display, "", no_update

        # Crawl all sites — each with its own session/cookies
        all_links = []
        search_query = " ".join(keywords[:8])  # Use top keywords as search query
        for site_url in url_list:
            stop_flag.check()
            site_domain = urlparse(normalize_url(site_url)).netloc.removeprefix("www.")

            # Detect search engines — search instead of crawl
            if is_search_engine(site_domain):
                logs.append(f"🔎 Searching {site_domain} for: {search_query[:50]}...")
                yield "\n".join(logs), kw_display, "", no_update
                try:
                    site_links = search_engine_crawl(site_domain, search_query, max_links=30)
                    logs.append(f"  📄 Found {len(site_links)} results from {site_domain}")
                    all_links.extend(site_links)
                except Exception as e:
                    logs.append(f"  ⚠️ Search failed: {e}")
                    _research_issues.append(f"Search {site_domain}: {e}")
            else:
                session = build_session_for_url(site_url, cookies_text)
                if session.cookies:
                    logs.append(f"🔐 Loaded {len(session.cookies)} cookies for {urlparse(normalize_url(site_url)).netloc}")
                depth_label = " (deep)" if crawl_deep else ""
                logs.append(f"🕷️ Crawling {site_url}{depth_label}...")
                yield "\n".join(logs), kw_display, "", no_update

                # Log callback for deep crawl progress
                def _log_update(msg, _logs=logs, _kw=kw_display):
                    _logs.append(msg)

                try:
                    if crawl_deep:
                        site_links = crawl_site_deep(
                            site_url, keywords, session=session,
                            max_links=150, max_sub_pages=10,
                            log_fn=_log_update,
                        )
                    else:
                        site_links = crawl_site_links(site_url, session=session, max_links=80)
                    logs.append(f"  📄 Found {len(site_links)} links from {urlparse(normalize_url(site_url)).netloc}")
                    all_links.extend(site_links)
                except Exception as e:
                    logs.append(f"  ⚠️ Failed: {e}")
                    _research_issues.append(f"Crawl {urlparse(normalize_url(site_url)).netloc}: {e}")

            yield "\n".join(logs), kw_display, "", no_update

        links = all_links
        logs.append(f"📄 Total: {len(links)} article links from {len(url_list)} site(s)")
        yield "\n".join(logs), kw_display, "", no_update

        # Phase 1: AI judges which titles are relevant (keyword pre-filtered internally)
        stop_flag.check()
        logs.append(f"🧠 AI filtering from {len(links)} titles...")
        yield "\n".join(logs), kw_display, "", no_update

        candidates = ai_filter_titles(links, prompt, ext_model, keywords=keywords)
        candidates.sort(key=lambda x: x["score"], reverse=True)

        if not candidates:
            logs.append("📊 AI found no relevant titles")
            yield "\n".join(logs), kw_display, "No matching articles found. Try a different topic or site.", gr.update(choices=[], value="")
            return

        stop_flag.check()
        logs.append(f"📊 AI selected {len(candidates)} relevant articles")
        for c in candidates:
            logs.append(f"  → [{int(c['score']*100)}%] {c['title'][:60]} — {c['reason']}")
        yield "\n".join(logs), kw_display, "", no_update

        # Phase 2: AI reads top candidates and scores relevance
        top = candidates[:max_articles]
        logs.append(f"📖 AI deep-reading top {len(top)} of {len(candidates)} articles...")
        yield "\n".join(logs), kw_display, "", no_update

        results = []
        for i, article in enumerate(top):
            stop_flag.check()
            logs.append(f"  📖 [{i+1}/{len(top)}] {article['title'][:50]}...")
            yield "\n".join(logs), kw_display, "", no_update

            art_session = build_session_for_url(article["url"], cookies_text)
            ai_result = ai_score_article(article["url"], prompt, art_session, ext_model)
            combined = (article["score"] * 0.4) + (ai_result["score"] * 0.6)
            reason = ai_result.get("reason") or article.get("reason", "")
            results.append({**article, "score": combined, "reason": reason})

        results.sort(key=lambda x: x["score"], reverse=True)
        global _research_results
        _research_results = results

        # Phase 3: AI cross-checks articles to group related topics
        stop_flag.check()
        if len(results) > 1:
            logs.append("🔗 AI cross-checking articles for related topics...")
            yield "\n".join(logs), kw_display, "", no_update

            try:
                titles_for_grouping = "\n".join(f"{i}. {r['title']}" for i, r in enumerate(results))
                raw = chat(get_prompts().group_articles(prompt, titles_for_grouping), ref_model)
                raw = raw.strip().replace("```json", "").replace("```", "").strip()
                topic_groups = json.loads(raw)
                logs.append(f"📊 Found {len(topic_groups)} topic groups:")
                for tg in topic_groups:
                    article_count = len(tg.get("articles", []))
                    logs.append(f"  • {tg['topic']} ({article_count} articles) — {tg.get('summary', '')[:60]}")
                # Tag each result with its topic group
                for tg in topic_groups:
                    for idx in tg.get("articles", []):
                        if 0 <= idx < len(results):
                            results[idx]["topic_group"] = tg["topic"]
            except Exception as e:
                logs.append(f"  ⚠️ Grouping failed: {e}")

            yield "\n".join(logs), kw_display, "", no_update

        # Format output + build dropdown choices
        output_lines = []
        url_choices = []
        current_topic = None
        for i, r in enumerate(results, 1):
            pct = int(r["score"] * 100)
            reason = f" — {r['reason']}" if r.get("reason") else ""
            topic = r.get("topic_group", "")
            if topic and topic != current_topic:
                output_lines.append(f"\n── {topic} ──")
                current_topic = topic
            output_lines.append(f"{i}. [{pct}% match] {r['title']}{reason}\n   {r['url']}")
            url_choices.append(f"{r['title'][:60]}  |  {r['url']}")

        result_text = "\n".join(output_lines)
        logs.append(f"\n✅ Found {len(results)} relevant articles")
        # Prefer a fetchable article as the default selection
        top_choice = ""
        for r in results:
            top_choice = f"{r['title'][:60]}  |  {r['url']}"
            break
        if not top_choice and url_choices:
            top_choice = url_choices[0]
        yield "\n".join(logs), kw_display, result_text, gr.update(choices=url_choices, value=top_choice)

    except GeneratorExit:
        logs.append("\n🛑 Stopped by user")
        yield "\n".join(logs), "", "", no_update
        return
    except Exception as e:
        logs.append(f"❌ Error: {e}")
        yield "\n".join(logs), "", "", no_update


def run_pipeline(url, lang1, lang2, num_slides, theme, open_keynote, model, cookies_text, user_prompt="", extra_urls=None, stop_flag=None):
    """Generate keynote from one or more article URLs. If extra_urls provided, combines content."""
    logs = []
    issues = list(_research_issues)  # Carry over any crawl issues
    lang2 = lang2 if lang2 and lang2 != "None" else ""
    if stop_flag is None:
        stop_flag = StopFlag()

    # Resolve refinement model for translation & slide generation
    available = get_available_models()
    ref_model = pick_stage_model("refinement", available, model)

    # Combine primary URL with any extra URLs
    all_urls = [url] + (extra_urls or [])
    # Resolve Google News URLs and deduplicate
    seen = set()
    unique_urls = []
    for u in all_urls:
        u = resolve_google_news_url(u)  # Decode Google News redirect
        if u in seen:
            continue
        seen.add(u)
        unique_urls.append(u)
    if not unique_urls:
        unique_urls = [url]

    try:
        lang_label = f"{lang1} + {lang2}" if lang2 else lang1
        if user_prompt:
            logs.append(f"📝 User prompt: {user_prompt[:80]}...")
        logs.append(f"🌍 Languages: {lang_label}")
        logs.append(f"🤖 Refinement model: {ref_model.split('/')[-1]}")
        if len(unique_urls) > 1:
            logs.append(f"📰 Combining {len(unique_urls)} related articles")

        # Fetch all article content
        stop_flag.check()
        all_content = []
        fetched_sources = []  # Track successfully fetched article URLs
        for i, article_url in enumerate(unique_urls):
            logs.append(f"📥 Fetching [{i+1}/{len(unique_urls)}] {urlparse(normalize_url(article_url)).path[:50]}...")
            yield "\n".join(logs), None, ""

            session = build_session_for_url(article_url, cookies_text)
            try:
                text = fetch_url(article_url, session)
                logs.append(f"  ✅ {len(text)} characters")
                all_content.append(text)
                fetched_sources.append(article_url)
            except Exception as e:
                issues.append(f"Fetch {article_url}: {e}")
                logs.append(f"  ⚠️ {e}")
                try:
                    text = requests.get(f"https://r.jina.ai/{article_url}", timeout=30).text
                    logs.append(f"  ✅ Fallback: {len(text)} characters")
                    all_content.append(text)
                    fetched_sources.append(article_url)
                except Exception as e2:
                    issues.append(f"Fallback {article_url}: {e2}")
                    logs.append(f"  ⚠️ Fallback also failed")
            yield "\n".join(logs), None, ""

        content = "\n\n---\n\n".join(all_content) if all_content else f"[No content fetched from {url}]"

        # Fetch 1 representative image from the primary article
        logs.append("🖼️ Looking for article image...")
        yield "\n".join(logs), None, ""
        article_image = None
        for article_url in unique_urls:
            try:
                session = build_session_for_url(article_url, cookies_text)
                imgs = fetch_article_images(article_url, session, max_images=1)
                if imgs:
                    article_image = imgs[0]
                    logs.append(f"  ✅ Found image from {urlparse(normalize_url(article_url)).netloc}")
                    break
            except Exception:
                continue
        if not article_image:
            logs.append("  ℹ️ No suitable image found")
        yield "\n".join(logs), None, ""

        stop_flag.check()
        logs.append("🌐 Translating & summarizing...")
        yield "\n".join(logs), None, ""

        try:
            summary = translate_content(content, lang1, ref_model, user_prompt=user_prompt, lang2=lang2)
            logs.append("✅ Translation done")
        except Exception as e:
            issues.append(f"Translation error: {e}")
            summary = content[:500]
            logs.append(f"⚠️ Translation failed, using raw content excerpt")
        yield "\n".join(logs), None, summary

        stop_flag.check()
        batch_info = f" (in batches of 4)" if num_slides > 5 else ""
        logs.append(f"🧠 Generating {num_slides} slides{batch_info}...")
        yield "\n".join(logs), None, summary

        try:
            data = generate_slide_structure(content, lang1, num_slides, ref_model, source_url=url, source_title=url, user_prompt=user_prompt, lang2=lang2)
            # Attach all fetched sources for slide-level attribution
            data["all_sources"] = fetched_sources
            logs.append(f"✅ Structure ready: {len(data['slides'])} slides")
        except Exception as e:
            issues.append(f"Slide generation error: {e}")
            logs.append(f"⚠️ Slide generation failed: {e}")
            yield "\n".join(logs), None, summary
            return
        yield "\n".join(logs), None, summary

        if issues:
            logs.append(f"⚠️ {len(issues)} issue(s) will be noted in the final slide")

        logs.append("📊 Building PPTX...")
        yield "\n".join(logs), None, summary

        data["summary"] = summary
        pptx_path = build_pptx(data, theme, issues=issues, cover_image=article_image)
        logs.append(f"✅ Saved: {os.path.basename(pptx_path)}")
        yield "\n".join(logs), None, summary

        # Export to PDF via LibreOffice (headless) and open in Keynote
        logs.append("📄 Exporting to PDF...")
        yield "\n".join(logs), None, summary

        pdf_path = pptx_path.replace(".pptx", ".pdf")
        download_files = []

        # PDF export via LibreOffice headless
        pdf_ok = False
        soffice_paths = [
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            "soffice",
            "libreoffice",
        ]
        for soffice in soffice_paths:
            try:
                subprocess.run(
                    [soffice, "--headless", "--convert-to", "pdf",
                     "--outdir", os.path.dirname(pptx_path), pptx_path],
                    timeout=120, check=True, capture_output=True,
                )
                if os.path.exists(pdf_path):
                    logs.append(f"✅ PDF: {os.path.basename(pdf_path)}")
                    download_files.append(pdf_path)
                    pdf_ok = True
                    break
            except (FileNotFoundError, subprocess.CalledProcessError):
                continue

        if not pdf_ok:
            logs.append("⚠️ PDF export unavailable (install LibreOffice for PDF)")

        download_files.append(pptx_path)
        logs.append(f"✅ PPTX: {os.path.basename(pptx_path)}")

        if open_keynote:
            if pdf_ok:
                subprocess.run(["open", pdf_path])
                logs.append(f"🎬 Opened {os.path.basename(pdf_path)}")
            else:
                subprocess.run(["open", "-a", "Keynote", pptx_path])
                logs.append(f"🎬 Opened in Keynote: {os.path.basename(pptx_path)}")

        yield "\n".join(logs), download_files, summary

    except GeneratorExit:
        logs.append("\n🛑 Stopped by user")
        yield "\n".join(logs), None, ""
        return
    except Exception as e:
        logs.append(f"❌ Error: {e}")
        yield "\n".join(logs), None, ""


def list_output_filenames() -> list[str]:
    """Return list of output filenames for dropdown."""
    if not os.path.isdir(OUTPUT_DIR):
        return []
    names = []
    for fname in sorted(os.listdir(OUTPUT_DIR), key=lambda f: os.path.getmtime(os.path.join(OUTPUT_DIR, f)), reverse=True):
        fpath = os.path.join(OUTPUT_DIR, fname)
        if os.path.isfile(fpath) and os.path.splitext(fname)[1].lower() in (".pptx", ".pdf"):
            names.append(fname)
    return names




# --- UI ---
css = """
/* === Web AI Tool — force dark mode === */
:root {
    color-scheme: dark !important;
}
*, *::before, *::after { box-sizing: border-box; }

html {
    background: #1a1a1a !important;
    color-scheme: dark !important;
}
body {
    font-family: 'Segoe UI', 'SF Pro Text', -apple-system, Roboto, sans-serif;
    background: #1a1a1a !important;
    color: #e0e0e0 !important;
}
/* Force dark on all Gradio wrappers */
.gradio-container,
.gradio-container *,
.main, .app, .contain,
div[class*="svelte"],
div[class*="gradio"] {
    color-scheme: dark !important;
}
/* Kill any white backgrounds from Gradio/Safari defaults */
.gradio-container div,
.gradio-container section,
.gradio-container form,
.gradio-container fieldset,
.gradio-container details,
.gradio-container summary {
    background-color: transparent;
}
.block, .gr-block, .gr-box, .gr-panel,
div[class*="block"], div[class*="panel"],
div[class*="group"], div[class*="form"] {
    background: transparent !important;
    border-color: #3a3a3a !important;
}

/* ── Container ── */
.gradio-container {
    max-width: 900px !important;
    margin: 0 auto !important;
    padding: 0 12px 12px !important;
    background: #2b2b2b;
    border-left: 1px solid #3a3a3a;
    border-right: 1px solid #3a3a3a;
    box-shadow: 0 0 40px rgba(0,0,0,0.5);
}

/* ── Kill Gradio's default spacing ── */
.gradio-container > .flex,
.gradio-container .form,
.gradio-container .contain,
.gradio-container .panel {
    gap: 5px !important;
    padding: 0 !important;
    margin: 0 !important;
}
.gradio-container .panel,
.gradio-container .gr-box {
    background: none !important;
    border: none !important;
    box-shadow: none !important;
}
.gradio-container .row,
.gradio-container .gr-row,
.gradio-container [class*="row"] { gap: 8px !important; align-items: end !important; }

/* ── Header bar ── */
.ext-header {
    background: linear-gradient(135deg, #1a73e8, #1558b0);
    padding: 12px 16px !important;
    margin: 0 -12px 10px !important;
    border-bottom: 1px solid #1558b0;
}
.ext-header h3, .ext-header p { margin: 0 !important; color: #fff !important; }
.ext-header h3 { font-size: 14px !important; font-weight: 600 !important; }
.ext-header p { font-size: 10.5px !important; opacity: 0.8; margin-top: 1px !important; }

/* ── Labels — plain text, no background ── */
.gradio-container label,
.gradio-container .label-wrap,
.gradio-container span[data-testid="block-info"],
.gradio-container span.text-gray-500,
.gradio-container .gradio-label,
.gradio-container label span {
    font-size: 10px !important;
    font-weight: 600 !important;
    color: #888 !important;
    text-transform: uppercase;
    letter-spacing: 0.5px;
    padding: 0 !important;
    margin: 0 0 1px !important;
    background: none !important;
    background-color: transparent !important;
    border: none !important;
    box-shadow: none !important;
    line-height: 1.4 !important;
}
/* Checkbox/radio labels — normal case, readable size */
.gradio-container label:has(input[type="checkbox"]) span,
.gradio-container label:has(input[type="radio"]) span {
    font-size: 12px !important;
    text-transform: none !important;
    letter-spacing: 0 !important;
    color: #ccc !important;
}

/* ── Inputs — consistent 34px height ── */
.gradio-container input[type="text"],
.gradio-container select,
.gradio-container .wrap input {
    background: #333 !important;
    border: 1px solid #444 !important;
    color: #e0e0e0 !important;
    border-radius: 5px !important;
    font-size: 12.5px !important;
    padding: 6px 8px !important;
    height: 34px !important;
    transition: border-color 0.15s;
}
.gradio-container .secondary-wrap {
    background: #333 !important;
    border: 1px solid #444 !important;
    color: #e0e0e0 !important;
    border-radius: 5px !important;
    font-size: 12.5px !important;
    padding: 4px 8px !important;
}
.gradio-container input[type="text"]:focus {
    border-color: #1a73e8 !important;
    outline: none !important;
    box-shadow: 0 0 0 2px rgba(26,115,232,0.2) !important;
}

/* ── Dropdown overrides ── */
.gradio-container .wrap.svelte-aqlk7e,
.gradio-container .wrap[data-testid] {
    background: #333 !important;
    border: 1px solid #444 !important;
    border-radius: 5px !important;
    min-height: 34px !important;
    max-height: 34px !important;
    padding: 0 8px !important;
    display: flex !important;
    align-items: center !important;
}
/* Multiselect dropdowns can grow */
.ext-site-tags .wrap.svelte-aqlk7e,
.ext-site-tags .wrap[data-testid] {
    max-height: 100px !important;
    min-height: 36px !important;
}

/* ── Tabs ── */
.ext-tabs > .tab-nav {
    background: #252525 !important;
    border-bottom: 1px solid #3a3a3a !important;
    gap: 0 !important;
    padding: 0 !important;
    margin: 0 -12px 6px !important;
}
.ext-tabs > .tab-nav button {
    background: transparent !important;
    color: #777 !important;
    border: none !important;
    border-bottom: 2px solid transparent !important;
    font-size: 11.5px !important;
    font-weight: 600 !important;
    padding: 8px 14px !important;
    margin: 0 !important;
    transition: all 0.15s;
}
.ext-tabs > .tab-nav button.selected {
    color: #e0e0e0 !important;
    border-bottom-color: #1a73e8 !important;
    background: rgba(26,115,232,0.08) !important;
}
.ext-tabs > .tab-nav button:hover:not(.selected) {
    color: #bbb !important;
    background: rgba(255,255,255,0.04) !important;
}
.ext-tabs .tabitem { padding: 6px 0 0 !important; }

/* ── Buttons ── */
.ext-generate-btn {
    background: linear-gradient(135deg, #1a73e8, #1558b0) !important;
    color: #fff !important;
    border: none !important;
    border-radius: 6px !important;
    font-size: 12.5px !important;
    font-weight: 600 !important;
    padding: 9px 0 !important;
    margin-top: 4px !important;
    cursor: pointer;
    transition: opacity 0.15s, transform 0.1s;
}
.ext-generate-btn:hover { opacity: 0.9 !important; }
.ext-generate-btn:active { transform: scale(0.98) !important; }


/* ── Slider compact ── */
.gradio-container input[type="range"] { accent-color: #1a73e8; }
.gradio-container input[type="number"] {
    background: #333 !important;
    border: 1px solid #444 !important;
    color: #ddd !important;
    border-radius: 5px !important;
    width: 60px !important;
    font-size: 12.5px !important;
    padding: 6px 8px !important;
    height: 34px !important;
}

/* ── Radio / checkbox — fully restore ── */
.gradio-container input[type="radio"],
.gradio-container input[type="checkbox"] {
    accent-color: #1a73e8;
    width: 16px !important;
    height: 16px !important;
    min-width: 16px !important;
    margin: 0 !important;
    padding: 0 !important;
    cursor: pointer !important;
    pointer-events: auto !important;
    position: static !important;
    opacity: 1 !important;
    flex-shrink: 0 !important;
    -webkit-appearance: auto !important;
    appearance: auto !important;
}
/* Every ancestor of a checkbox/radio must allow clicks and be visible */
.gradio-container *:has(input[type="checkbox"]),
.gradio-container *:has(input[type="radio"]) {
    pointer-events: auto !important;
    overflow: visible !important;
    max-height: none !important;
    height: auto !important;
    min-height: 0 !important;
}
.gradio-container label:has(input[type="checkbox"]),
.gradio-container label:has(input[type="radio"]) {
    display: flex !important;
    align-items: center !important;
    gap: 6px !important;
    padding: 4px 8px !important;
    cursor: pointer !important;
}
/* Radio/checkbox group container — align height with inputs */
.gradio-container .wrap:has(> label > input[type="radio"]),
.gradio-container .wrap:has(> label > input[type="checkbox"]) {
    min-height: 34px !important;
    max-height: none !important;
    display: flex !important;
    flex-wrap: wrap !important;
    align-items: center !important;
    gap: 0 !important;
}

/* ── Section divider ── */
.ext-section {
    border-top: 1px solid #333;
    margin: 6px 0 4px !important;
    padding: 0 !important;
}

/* ── Progress log — minimal overrides, preserve Gradio streaming flash ── */
@keyframes border-pulse {
    0%   { box-shadow: 0 0 0 0 rgba(139,195,74,0); }
    50%  { box-shadow: 0 0 10px 2px rgba(139,195,74,0.35); }
    100% { box-shadow: 0 0 0 0 rgba(139,195,74,0); }
}
.ext-progress textarea {
    font-family: 'SF Mono', 'Cascadia Code', 'Fira Code', monospace !important;
    font-size: 11px !important;
    line-height: 1.5 !important;
}
.ext-progress.ext-active {
    animation: border-pulse 1.5s ease-in-out infinite !important;
}

/* ── Summary / results ── */
.ext-summary textarea {
    font-size: 11.5px !important;
    line-height: 1.5 !important;
}

/* ── File download ── */
.gradio-container .file-preview,
.gradio-container .upload-button {
    background: #333 !important;
    border: 1px dashed #444 !important;
    border-radius: 5px !important;
    color: #aaa !important;
    font-size: 11px !important;
}

/* ── Scrollbar ── */
::-webkit-scrollbar { width: 5px; }
::-webkit-scrollbar-track { background: transparent; }
::-webkit-scrollbar-thumb { background: #444; border-radius: 3px; }
::-webkit-scrollbar-thumb:hover { background: #666; }

/* ── Accordion (cookies) ── */
.ext-accordion {
    border: 1px solid #3a3a3a !important;
    border-radius: 6px !important;
    background: #252525 !important;
    margin: 4px 0 !important;
    padding: 0 !important;
}
.ext-accordion > .label-wrap {
    padding: 6px 10px !important;
    background: none !important;
    color: #999 !important;
    font-size: 11px !important;
}
.ext-accordion > .label-wrap span { color: #999 !important; font-size: 11px !important; }
.ext-accordion .prose { padding: 0 !important; }
.ext-accordion > div:not(.label-wrap) {
    padding: 4px 10px 8px !important;
    gap: 4px !important;
}
.ext-accordion > div:not(.label-wrap) > div { gap: 4px !important; padding: 0 !important; margin: 0 !important; }
.ext-accordion > div:not(.label-wrap) > div > div { gap: 2px !important; padding: 0 !important; margin: 0 !important; }
.ext-accordion .row { gap: 6px !important; }
.ext-accordion button { margin: 0 !important; padding: 7px 0 !important; }
.ext-accordion label {
    background: none !important; border: none !important; box-shadow: none !important;
    padding: 0 !important; margin: 0 0 1px !important;
}

/* ── Site URL tags ── */
.ext-site-tags .token {
    background: #1a73e8 !important;
    color: #fff !important;
    border-radius: 12px !important;
    padding: 2px 10px !important;
    font-size: 11px !important;
    margin: 2px !important;
}
.ext-site-tags .token-remove {
    color: #fff !important;
    margin-left: 4px !important;
    cursor: pointer !important;
}
.ext-site-tags .secondary-wrap {
    min-height: 36px !important;
    flex-wrap: wrap !important;
    gap: 3px !important;
    padding: 4px 8px !important;
}
/* Hide the clear-all × and dropdown arrow */
.ext-site-tags .icon-wrap,
.ext-site-tags button[aria-label="Clear"],
.ext-site-tags .clear-btn,
.ext-site-tags > div > div > button:not(.token-remove),
.ext-site-tags svg.icon-clear {
    display: none !important;
}

/* ── Output files panel — scrollable container ── */
/* ── Output files panel ── */
.ext-output-panel {
    border: 1px solid #3a3a3a !important;
    background: #252525 !important;
    border-radius: 8px !important;
    margin: 8px 0 !important;
    padding: 0 !important;
    max-height: 400px !important;
    overflow-y: auto !important;
    overflow-x: hidden !important;
}
.ext-output-header {
    font-size: 12px !important;
    font-weight: 600 !important;
    color: #aaa !important;
    text-transform: uppercase !important;
    letter-spacing: 0.5px !important;
    padding: 10px 4px 6px !important;
    margin: 0 !important;
}
/* Confirm bar hidden by default */
#output-confirm-bar { display: none !important; }
#output-confirm-bar.ext-visible { display: flex !important; }

.ext-action-bar {
    padding: 6px 0 !important;
    margin: 0 !important;
    gap: 6px !important;
    flex-direction: row !important;
    justify-content: center !important;
    align-items: center !important;
}
.ext-action-bar button {
    font-size: 12px !important;
    font-weight: 500 !important;
    padding: 7px 16px !important;
    border-radius: 6px !important;
    min-height: 34px !important;
    margin: 0 !important;
    border: 1px solid #444 !important;
    background: #333 !important;
    color: #ddd !important;
    cursor: pointer !important;
    transition: background 0.15s !important;
}
.ext-action-bar button:hover {
    background: #3e3e3e !important;
}
.ext-action-bar button:last-child {
    border-color: #633 !important;
    color: #e57373 !important;
}
.ext-action-bar button:last-child:hover {
    background: #3a2020 !important;
}

/* ── Hide footer ── */
footer { display: none !important; }

/* ══════════════════════════════════════════════
   Mobile SPA — (≤768px)
   ══════════════════════════════════════════════ */
@media (max-width: 768px) {
    /* ── SPA shell — no horizontal scroll ── */
    html, body {
        overflow-x: hidden !important;
        -webkit-text-size-adjust: 100% !important;
    }
    body {
        background: #1e1e1e !important;
        -webkit-tap-highlight-color: transparent;
    }

    /* ── Nuclear: force ALL horizontal flex to stack ── */
    .gradio-container div[style*="flex-direction"],
    .gradio-container .flex,
    .gradio-container .gap {
        flex-direction: column !important;
        align-items: stretch !important;
    }
    .gradio-container .flex > *,
    .gradio-container .gap > * {
        width: 100% !important;
        min-width: 100% !important;
        flex: 1 1 100% !important;
    }
    /* Exception: keep radio/checkbox groups horizontal */
    .gradio-container .wrap:has(> label > input[type="radio"]),
    .gradio-container .wrap:has(> label > input[type="checkbox"]) {
        flex-direction: row !important;
        flex-wrap: wrap !important;
        align-items: center !important;
    }
    .gradio-container .wrap:has(> label > input[type="radio"]) > *,
    .gradio-container .wrap:has(> label > input[type="checkbox"]) > * {
        width: auto !important;
        min-width: auto !important;
        flex: 0 0 auto !important;
    }
    /* Exception: keep tab nav horizontal */
    .ext-tabs > .tab-nav {
        flex-direction: row !important;
    }
    .ext-tabs > .tab-nav > * {
        width: auto !important;
        min-width: auto !important;
        flex: 0 0 auto !important;
    }
    /* Exception: keep action bar horizontal */
    .ext-action-bar {
        flex-direction: row !important;
    }
    .ext-action-bar > * {
        width: auto !important;
        min-width: auto !important;
        flex: 1 1 0 !important;
    }

    /* ── Container ── */
    .gradio-container {
        max-width: 100% !important;
        padding: 0 8px 80px !important;
        margin: 0 !important;
        border: none !important;
        box-shadow: none !important;
        background: #1e1e1e !important;
    }

    /* ── Sticky header ── */
    .ext-header {
        position: sticky !important;
        top: 0 !important;
        z-index: 100 !important;
        margin: 0 -8px 8px !important;
        padding: 12px 14px !important;
        backdrop-filter: blur(10px) !important;
        -webkit-backdrop-filter: blur(10px) !important;
        background: linear-gradient(135deg, rgba(26,115,232,0.95), rgba(21,88,176,0.95)) !important;
    }
    .ext-header h3 { font-size: 16px !important; }
    .ext-header p { font-size: 11px !important; }

    /* ── Rows — force stack vertically on mobile ── */
    .gradio-container .row,
    .gradio-container .gr-row,
    .gradio-container [class*="row"],
    .gradio-container .flex-row,
    .gradio-container div[style*="flex-direction: row"],
    .gradio-container .form > div:has(> div + div) {
        flex-direction: column !important;
        flex-wrap: wrap !important;
        gap: 6px !important;
        align-items: stretch !important;
    }
    .gradio-container .row > *,
    .gradio-container .gr-row > *,
    .gradio-container [class*="row"] > * {
        min-width: 100% !important;
        max-width: 100% !important;
        flex: 1 1 100% !important;
        width: 100% !important;
    }

    /* ── Columns — stack ── */
    .gradio-container .column,
    .gradio-container .gr-column,
    .gradio-container [class*="column"] {
        min-width: 100% !important;
        width: 100% !important;
    }

    /* ── Tabs — horizontal scroll, pill style ── */
    .ext-tabs > .tab-nav {
        overflow-x: auto !important;
        margin: 0 -8px 8px !important;
        padding: 0 8px !important;
        -webkit-overflow-scrolling: touch;
        scrollbar-width: none;
        background: #252525 !important;
        position: sticky !important;
        top: 44px !important;
        z-index: 99 !important;
    }
    .ext-tabs > .tab-nav::-webkit-scrollbar { display: none; }
    .ext-tabs > .tab-nav button {
        padding: 10px 14px !important;
        font-size: 12px !important;
        white-space: nowrap !important;
        flex-shrink: 0 !important;
    }

    /* ── Buttons — 44px min touch target (Apple HIG) ── */
    .ext-generate-btn {
        padding: 14px 0 !important;
        font-size: 15px !important;
        margin-top: 8px !important;
        border-radius: 10px !important;
    }
    .ext-refresh-btn {
        height: 44px !important;
        min-width: 44px !important;
        font-size: 18px !important;
        border-radius: 8px !important;
    }
    button {
        min-height: 44px !important;
    }

    /* ── Inputs — 44px for touch ── */
    .gradio-container input[type="text"],
    .gradio-container select,
    .gradio-container .wrap input {
        height: 44px !important;
        font-size: 16px !important;  /* prevents iOS zoom on focus */
        padding: 8px 12px !important;
        border-radius: 8px !important;
    }
    .gradio-container input[type="number"] {
        height: 44px !important;
        width: 100% !important;
        font-size: 16px !important;
        border-radius: 8px !important;
    }

    /* ── Textareas ── */
    .gradio-container textarea {
        font-size: 14px !important;
        border-radius: 8px !important;
    }

    /* ── Dropdowns — taller ── */
    .gradio-container .wrap.svelte-aqlk7e,
    .gradio-container .wrap[data-testid] {
        min-height: 44px !important;
        max-height: 44px !important;
        border-radius: 8px !important;
    }
    .ext-site-tags .wrap.svelte-aqlk7e,
    .ext-site-tags .wrap[data-testid] {
        max-height: 140px !important;
        min-height: 48px !important;
    }

    /* ── Radio/checkbox — larger touch ── */
    .gradio-container input[type="radio"],
    .gradio-container input[type="checkbox"] {
        width: 22px !important;
        height: 22px !important;
        min-width: 22px !important;
    }
    .gradio-container label:has(input[type="checkbox"]),
    .gradio-container label:has(input[type="radio"]) {
        padding: 8px 12px !important;
        gap: 10px !important;
    }
    .gradio-container label:has(input[type="checkbox"]) span,
    .gradio-container label:has(input[type="radio"]) span {
        font-size: 14px !important;
    }

    /* ── Labels ── */
    .gradio-container label,
    .gradio-container .label-wrap,
    .gradio-container .gradio-label,
    .gradio-container label span {
        font-size: 11px !important;
    }

    /* ── Accordion — card style ── */
    .ext-accordion {
        border-radius: 10px !important;
        margin: 6px 0 !important;
    }
    .ext-accordion > .label-wrap {
        padding: 12px 14px !important;
    }
    .ext-accordion > .label-wrap span {
        font-size: 13px !important;
    }
    .ext-accordion > div:not(.label-wrap) {
        padding: 8px 12px 12px !important;
    }

    /* ── Site tags ── */
    .ext-site-tags .token {
        padding: 4px 10px !important;
        font-size: 11px !important;
        margin: 2px !important;
        border-radius: 14px !important;
        max-width: 85vw !important;
        overflow: hidden !important;
        text-overflow: ellipsis !important;
    }

    /* ── Progress/summary ── */
    .ext-progress textarea,
    .ext-summary textarea {
        font-size: 12px !important;
        border-radius: 8px !important;
    }

    /* ── File download ── */
    .gradio-container a[download] {
        font-size: 14px !important;
        padding: 8px 0 !important;
    }

    /* ── Section divider ── */
    .ext-section {
        margin: 10px 0 8px !important;
    }

    /* ── Hide Gradio branding/extra chrome ── */
    .built-with { display: none !important; }
    .gradio-container > .flex > .flex:empty { display: none !important; }
}

/* ── Extra small (≤480px) ── */
@media (max-width: 480px) {
    .ext-header h3 { font-size: 15px !important; }
    .ext-header p { font-size: 10px !important; }

    .ext-tabs > .tab-nav button {
        padding: 8px 10px !important;
        font-size: 11px !important;
    }

    .gradio-container label,
    .gradio-container label span {
        font-size: 10px !important;
    }

    .gradio-container {
        padding: 0 1px 80px !important;
    }
}
"""

# Load saved provider preference
_startup_prefs = load_saved_prefs()
_saved_provider = _startup_prefs.get("ai_provider", "LM Studio")
_saved_custom_url = _startup_prefs.get("ai_provider_url", "")
if _saved_provider in AI_PROVIDERS:
    set_provider_base(AI_PROVIDERS[_saved_provider][0])
elif _saved_custom_url:
    set_provider_base(_saved_custom_url)

available_models = get_available_models()
_has_ext = any(EXTRACTION_MODEL_PATTERN.lower() in m.lower() for m in available_models)
_has_ref = any(REFINEMENT_MODEL_PATTERN.lower() in m.lower() for m in available_models)
if _has_ext and _has_ref:
    available_models = ["auto"] + available_models
    default_model = "auto"
else:
    default_model = pick_default_model(available_models)

progress_js = """
() => {
    // Inject viewport meta for mobile
    if (!document.querySelector('meta[name="viewport"]')) {
        const vp = document.createElement('meta');
        vp.name = 'viewport';
        vp.content = 'width=device-width, initial-scale=1.0, maximum-scale=5.0, user-scalable=yes';
        document.head.appendChild(vp);
    }
    // PWA meta tags
    ['apple-mobile-web-app-capable', 'mobile-web-app-capable'].forEach(n => {
        if (!document.querySelector(`meta[name="${n}"]`)) {
            const m = document.createElement('meta');
            m.name = n; m.content = 'yes';
            document.head.appendChild(m);
        }
    });
    if (!document.querySelector('meta[name="theme-color"]')) {
        const tc = document.createElement('meta');
        tc.name = 'theme-color'; tc.content = '#1a1a1a';
        document.head.appendChild(tc);
    }
    // Force dark mode for Safari — inject early to prevent white flash
    if (!document.querySelector('#force-dark-mode')) {
        const ds = document.createElement('style');
        ds.id = 'force-dark-mode';
        ds.textContent = `
            :root { color-scheme: dark !important; }
            html, body { background: #1a1a1a !important; color: #e0e0e0 !important; }
            .block, .gr-block, .gr-box, .gr-panel,
            div[class*="block"], div[class*="panel"],
            div[class*="group"], div[class*="form"],
            .gradio-container div { background-color: transparent; }
        `;
        document.head.insertBefore(ds, document.head.firstChild);
    }

    // Mobile: strip Gradio wrapper padding and lock width
    if (/Mobi|Android|iPhone|iPad/i.test(navigator.userAgent)) {
        const lockStyle = document.createElement('style');
        lockStyle.textContent = `
            body > div, body > div > div, body > div > div > div {
                padding: 0 !important;
                margin: 0 !important;
                overflow-x: hidden !important;
            }
            table { display: block; overflow-x: auto; }
        `;
        document.head.appendChild(lockStyle);
    }


    // Fix Gradio's orphaned label[for] attributes (suppresses browser warning)
    function fixOrphanedLabels() {
        document.querySelectorAll('label[for]').forEach(label => {
            if (!document.getElementById(label.getAttribute('for'))) {
                label.removeAttribute('for');
            }
        });
    }
    fixOrphanedLabels();

    const observer = new MutationObserver(() => {
        fixOrphanedLabels();
        // Progress bar pulse
        document.querySelectorAll('.ext-progress textarea').forEach(ta => {
            const wrapper = ta.closest('.ext-progress');
            if (!wrapper) return;
            if (ta.value && ta.value.trim()) {
                // Auto-scroll to bottom
                ta.scrollTop = ta.scrollHeight;
                wrapper.classList.add('ext-active');
                clearTimeout(wrapper._pulseTimer);
                wrapper._pulseTimer = setTimeout(() => {
                    wrapper.classList.remove('ext-active');
                }, 3000);
            }
        });
        // Auto-download: prefer PDF, skip .key and .pptx
        document.querySelectorAll('a[download]').forEach(a => {
            if (a.dataset.autoDownloaded) return;
            if (a.href && a.href.includes('/file=')) {
                a.dataset.autoDownloaded = 'true';
                if (a.href.endsWith('.pdf') || a.download.endsWith('.pdf')) {
                    setTimeout(() => a.click(), 500);
                }
            }
        });
    });
    setTimeout(() => {
        observer.observe(document.body, { childList: true, subtree: true, characterData: true });
    }, 1000);
}
"""

mobile_head = """
<meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0, user-scalable=no">
<meta name="apple-mobile-web-app-capable" content="yes">
<meta name="apple-mobile-web-app-status-bar-style" content="black-translucent">
<meta name="mobile-web-app-capable" content="yes">
<meta name="theme-color" content="#1a1a1a">
<link rel="apple-touch-icon" href="data:image/svg+xml,<svg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 100 100'><text y='.9em' font-size='90'>🌐</text></svg>">
"""

with gr.Blocks(title="Web AI Tool") as app:
    gr.HTML("""
        <div class="ext-header">
            <h3>🌐 Web AI Tool</h3>
            <p>Research · Translate · Export to Keynote</p>
        </div>
    """)

    # Shared controls at the top
    with gr.Row():
        provider_input = gr.Dropdown(
            choices=list(AI_PROVIDERS.keys()),
            value=_saved_provider,
            label="AI Provider",
            scale=2,
        )
        provider_url = gr.Textbox(
            value=AI_PROVIDERS.get(_saved_provider, ("http://localhost:1234",))[0] if _saved_provider != "Custom" else _saved_custom_url,
            label="Endpoint URL",
            scale=3,
            interactive=(_saved_provider == "Custom"),
        )
    model_input = gr.Dropdown(
        choices=available_models,
        value=default_model,
        label="Model",
        allow_custom_value=True,
    )

    categories = load_categories()
    cat_names = list(categories.keys())
    saved_sites = load_saved_sites()
    sites_history = load_sites_history()
    all_suggested = get_all_suggested_sites(categories)
    all_choices = list(dict.fromkeys(all_suggested + sites_history + saved_sites))
    prefs = load_saved_prefs()
    default_cat = prefs.get("last_category", cat_names[0] if cat_names else "Finance")
    if default_cat not in cat_names:
        default_cat = cat_names[0] if cat_names else "Finance"

    with gr.Row():
        site_category = gr.Radio(
            choices=cat_names,
            value=default_cat,
            label="Category",
            scale=3,
        )
        site_category_all = gr.Checkbox(value=False, label="All categories", scale=1)

    research_site = gr.Dropdown(
        label="Site URLs (type and press Enter to add)",
        choices=all_choices,
        value=categories.get(default_cat, []),
        multiselect=True,
        allow_custom_value=True,
        elem_classes=["ext-site-tags"],
    )

    with gr.Accordion("⚙️ Manage Categories", open=False, elem_classes=["ext-accordion"]):
        with gr.Row():
            cat_edit_select = gr.Dropdown(choices=cat_names, value=default_cat, label="Edit Category", scale=3)
            cat_new_name = gr.Textbox(label="New name", placeholder="Rename or new category", scale=2)
            cat_add_btn = gr.Button("Add / Rename", scale=1, min_width=80)
            cat_del_btn = gr.Button("Delete", variant="stop", scale=1, min_width=60)
        cat_edit_sites = gr.Dropdown(
            label="Sites in category (edit, add, remove)",
            choices=all_choices,
            value=categories.get(default_cat, []),
            multiselect=True,
            allow_custom_value=True,
        )
        cat_save_btn = gr.Button("Save Category", variant="primary")

    cat_status = gr.HTML("")

    with gr.Accordion("🔐 Session / Cookies", open=False, elem_classes=["ext-accordion"]):
        with gr.Row():
            cookie_browser = gr.Dropdown(
                choices=["Chrome", "Firefox", "Safari", "Edge"],
                value="Chrome",
                label="Browser",
                scale=3,
            )
            cookie_load_btn = gr.Button("Load Cookies", elem_classes=["ext-generate-btn"], scale=4)
        cookies_input = gr.Textbox(
            label="Cookies",
            placeholder="Auto-filled from Site URL, or paste manually (name=val; name2=val2)",
            lines=2,
        )
        cookie_status = gr.HTML('<p style="font-size:10px;color:#666;margin:2px 0 0;">Loads cookies for the domain in Site URL above.</p>')

    with gr.Tabs(elem_classes=["ext-tabs"]):

        # ── Research Tab ──
        with gr.TabItem("🔍 Research"):
            with gr.Row():
                with gr.Column(scale=1):
                    research_prompt = gr.Textbox(
                        label="Topic / Prompt",
                        placeholder="e.g. Impact of AI regulation on tech companies",
                        lines=3,
                    )
                    with gr.Row():
                        research_max = gr.Number(value=prefs["max_articles"], label="Articles to deep-scan", precision=0, minimum=1, scale=2)
                        crawl_depth = gr.Radio(["Standard", "Deep"], value=prefs.get("crawl_depth", "Standard"), label="Crawl Depth", scale=2)
                    auto_keynote = gr.Checkbox(value=prefs["auto_keynote"], label="Auto-generate Keynote from top result")
                    with gr.Row():
                        research_btn = gr.Button("Search Articles", variant="primary", elem_classes=["ext-generate-btn"], scale=3)
                        research_stop = gr.Button("Stop", variant="stop", interactive=False, scale=1)
            research_log = gr.Textbox(label="Progress", lines=8, interactive=False, elem_classes=["ext-progress"])
            research_kw  = gr.Textbox(label="Extracted Keywords", lines=1, interactive=False)
            research_out = gr.Textbox(label="Matching Articles", lines=10, interactive=False, elem_classes=["ext-summary"])

            gr.HTML('<div class="ext-section"></div>')
            article_picker = gr.Dropdown(label="Select article for Keynote", choices=[], interactive=True)
            with gr.Row():
                pick_lang1  = gr.Dropdown(["zh-TW", "English", "Japanese", "Korean", "Spanish", "French", "German"], value=prefs["lang1"], label="Primary Lang", scale=2)
                pick_lang2  = gr.Dropdown(["English", "zh-TW", "Japanese", "Korean", "Spanish", "French", "German", "None"], value=prefs["lang2"], label="Secondary Lang", scale=2)
                pick_slides = gr.Number(value=prefs["slides"], label="Slides", precision=0, minimum=1, scale=1)
                pick_theme  = gr.Radio(["Dark", "Light", "Blue"], value=prefs["theme"], label="Theme", scale=2)
            with gr.Row():
                generate_from_research = gr.Button("Generate Keynote from Article", variant="primary", elem_classes=["ext-generate-btn"], scale=3)
                pick_stop = gr.Button("Stop", variant="stop", interactive=False, scale=1)
            pick_log     = gr.Textbox(label="Keynote Progress", lines=5, interactive=False, elem_classes=["ext-progress"])
            pick_file    = gr.File(label="Download")
            pick_summary = gr.Textbox(label="Summary", lines=4, interactive=False, elem_classes=["ext-summary"])

        # ── Keynote Tab ──
        with gr.TabItem("📊 Keynote"):
            with gr.Row():
                with gr.Column(scale=1):
                    url_input = gr.Textbox(label="URL", placeholder="https://example.com/article", lines=1)
                    with gr.Row():
                        lang1      = gr.Dropdown(["zh-TW", "English", "Japanese", "Korean", "Spanish", "French", "German"], value=prefs["lang1"], label="Primary Lang", scale=2)
                        lang2      = gr.Dropdown(["English", "zh-TW", "Japanese", "Korean", "Spanish", "French", "German", "None"], value=prefs["lang2"], label="Secondary Lang", scale=2)
                        num_slides = gr.Number(value=prefs["slides"], label="Slides", precision=0, minimum=1, scale=1)
                    with gr.Row():
                        theme        = gr.Radio(["Dark", "Light", "Blue"], value=prefs["theme"], label="Theme", scale=2)
                        open_keynote = gr.Checkbox(value=True, label="Auto-open file", scale=1)
                    with gr.Row():
                        btn = gr.Button("Generate Keynote", variant="primary", elem_classes=["ext-generate-btn"], scale=3)
                        keynote_stop = gr.Button("Stop", variant="stop", interactive=False, scale=1)
                with gr.Column(scale=1):
                    log_output     = gr.Textbox(label="Progress", lines=6, interactive=False, elem_classes=["ext-progress"])
                    file_output    = gr.File(label="Download")
                    summary_output = gr.Textbox(label="Summary", lines=5, interactive=False, elem_classes=["ext-summary"])

    # ── Output Files ──
    gr.HTML('<div class="ext-section"></div>')
    with gr.Group(elem_classes=["ext-output-panel"]):
        gr.HTML('<div class="ext-output-header">📁 Output Files</div>')
        output_file_picker = gr.Dropdown(
            label="Select file",
            choices=list_output_filenames(),
            interactive=True,
        )
        with gr.Row(elem_classes=["ext-action-bar"]):
            output_refresh_btn = gr.Button("🔄 Refresh", scale=1)
            output_download_btn = gr.Button("📥 Download", scale=1)
            output_open_btn = gr.Button("📂 Open on Mac", scale=1)
            output_delete_btn = gr.Button("🗑️ Delete", variant="stop", scale=1)
        output_status = gr.HTML("")
        with gr.Row(elem_classes=["ext-action-bar"], visible=False) as output_confirm_row:
            output_confirm_btn = gr.Button("⚠️ Confirm Delete", variant="stop", scale=1)
            output_cancel_btn = gr.Button("Cancel", scale=1)
        output_download = gr.File(label="Download", visible=False)

    # ── Live sync timer — polls shared progress so all clients stay updated ──
    sync_timer = gr.Timer(2)  # poll every 2 seconds
    _last_poll_version = gr.State(0)

    def _poll_progress(last_ver):
        """Return current progress if changed, otherwise skip update."""
        state, ver = _progress.snapshot()
        if ver == last_ver:
            # No change — return no_update for everything to avoid flicker
            return (gr.update(), gr.update(), gr.update(),
                    gr.update(), gr.update(), ver)
        running = state.get("running")
        # Only push fields relevant to the active pipeline
        r_log = gr.update(value=state["research_log"]) if state["research_log"] else gr.update()
        r_kw = gr.update(value=state["research_kw"]) if state["research_kw"] else gr.update()
        r_out = gr.update(value=state["research_out"]) if state["research_out"] else gr.update()
        p_log = gr.update(value=state["pick_log"]) if state["pick_log"] else gr.update()
        k_log = gr.update(value=state["keynote_log"]) if state["keynote_log"] else gr.update()
        return r_log, r_kw, r_out, p_log, k_log, ver

    sync_timer.tick(
        _poll_progress,
        inputs=[_last_poll_version],
        outputs=[research_log, research_kw, research_out, pick_log, log_output, _last_poll_version],
    )

    # ── Model auto-refresh timer (bound after handlers defined) ──
    model_timer = gr.Timer(30)

    # ── Event handlers ──
    def on_provider_change(provider_name):
        """Switch AI provider — update URL, save preference, refresh models."""
        if provider_name in AI_PROVIDERS:
            base_url = AI_PROVIDERS[provider_name][0]
            is_custom = provider_name == "Custom"
            set_provider_base(base_url)
            save_prefs(ai_provider=provider_name, ai_provider_url=base_url)
            models = get_available_models()
            has_ext = any(EXTRACTION_MODEL_PATTERN.lower() in m.lower() for m in models)
            has_ref = any(REFINEMENT_MODEL_PATTERN.lower() in m.lower() for m in models)
            choices = (["auto"] if has_ext and has_ref else []) + models
            default = "auto" if has_ext and has_ref else pick_default_model(models)
            status = f"{len(models)} models" if models else "not connected"
            return (
                gr.update(value=base_url, interactive=is_custom),
                gr.update(choices=choices, value=default, label=f"Model ({status})"),
            )
        return gr.update(), gr.update()

    def on_provider_url_change(url):
        """Apply custom endpoint URL."""
        url = (url or "").strip().rstrip("/")
        if not url:
            return gr.update()
        set_provider_base(url)
        save_prefs(ai_provider_url=url)
        models = get_available_models()
        has_ext = any(EXTRACTION_MODEL_PATTERN.lower() in m.lower() for m in models)
        has_ref = any(REFINEMENT_MODEL_PATTERN.lower() in m.lower() for m in models)
        choices = (["auto"] if has_ext and has_ref else []) + models
        default = "auto" if has_ext and has_ref else pick_default_model(models)
        status = f"{len(models)} models" if models else "not connected"
        return gr.update(choices=choices, value=default, label=f"Model ({status})")

    _last_known_models = {"list": []}

    def refresh_models():
        """Full refresh — resets selection to best default."""
        models = get_available_models()
        has_ext = any(EXTRACTION_MODEL_PATTERN.lower() in m.lower() for m in models)
        has_ref = any(REFINEMENT_MODEL_PATTERN.lower() in m.lower() for m in models)
        choices = (["auto"] if has_ext and has_ref else []) + models
        default = "auto" if has_ext and has_ref else pick_default_model(models)
        status = f"{len(models)} models" if models else "not connected"
        _last_known_models["list"] = choices
        return gr.update(choices=choices, value=default, label=f"Model ({status})")

    def poll_models(current_model):
        """Periodic check — only update choices if models changed, preserve selection."""
        models = get_available_models()
        has_ext = any(EXTRACTION_MODEL_PATTERN.lower() in m.lower() for m in models)
        has_ref = any(REFINEMENT_MODEL_PATTERN.lower() in m.lower() for m in models)
        choices = (["auto"] if has_ext and has_ref else []) + models
        if choices == _last_known_models["list"]:
            return gr.update()  # no change
        _last_known_models["list"] = choices
        status = f"{len(models)} models" if models else "not connected"
        # Keep current selection if still valid
        value = current_model if current_model in choices else (choices[0] if choices else "")
        return gr.update(choices=choices, value=value, label=f"Model ({status})")

    def on_load_cookies(browser, site_urls):
        # Support both list (multiselect) and string
        if isinstance(site_urls, list):
            url_list = [u.strip() for u in site_urls if u and u.strip()]
        else:
            url_list = [u.strip() for u in (site_urls or "").split(";") if u.strip()]
        if not url_list:
            return "", '<p style="font-size:10px;color:#e74c3c;">Add at least one Site URL first.</p>'
        all_cookies = []
        domains_loaded = []
        for raw_url in url_list:
            parsed = urlparse(raw_url if "://" in raw_url else f"https://{raw_url}")
            domain = parsed.netloc.removeprefix("www.")
            if not domain:
                continue
            result = load_browser_cookies(browser, domain)
            if not result.startswith("Error:"):
                all_cookies.append(result)
                domains_loaded.append(domain)
        if not all_cookies:
            return "", '<p style="font-size:10px;color:#e74c3c;">No cookies found for any domain.</p>'
        combined = "; ".join(all_cookies)
        count = combined.count("=")
        domains_str = ", ".join(domains_loaded)
        return combined, f'<p style="font-size:10px;color:#8bc34a;">Loaded {count} cookies for {domains_str} from {browser}</p>'

    def on_sites_change(sites):
        """Save current site selections for the research pipeline."""
        if isinstance(sites, list):
            save_sites(sites)

    def on_category_change(cat_name, use_all):
        """Switch sites dropdown to show the selected category's sites."""
        cats = load_categories()
        if use_all:
            all_sites = []
            seen = set()
            for sites in cats.values():
                for s in sites:
                    if s not in seen:
                        all_sites.append(s)
                        seen.add(s)
            return gr.update(value=all_sites)
        sites = cats.get(cat_name, [])
        save_prefs(last_category=cat_name)
        return gr.update(value=sites)

    def on_all_categories_toggle(use_all, cat_name):
        """Toggle between all categories and single category."""
        return on_category_change(cat_name, use_all)

    def on_cat_edit_select(cat_name):
        """Load category sites into the edit dropdown."""
        cats = load_categories()
        return gr.update(value=cats.get(cat_name, []))

    def _cat_status_html(msg, color="#8bc34a"):
        return f'<p style="font-size:11px;color:{color};margin:4px 0;">{msg}</p>'

    def on_cat_add_rename(edit_cat, new_name, edit_sites):
        """Add a new category or rename an existing one."""
        cats = load_categories()
        new_name = (new_name or "").strip()
        if not new_name:
            # Just save sites to the selected category
            if edit_cat:
                cats[edit_cat] = [s.strip() for s in (edit_sites or []) if s and s.strip()]
                save_categories(cats)
            cat_names = list(cats.keys())
            msg = _cat_status_html(f"✅ Saved {len(cats.get(edit_cat, []))} sites to \"{edit_cat}\"")
            return (gr.update(choices=cat_names),
                    gr.update(choices=cat_names, value=edit_cat),
                    gr.update(value=cats.get(edit_cat, [])),
                    msg)

        if edit_cat and edit_cat in cats and new_name != edit_cat:
            # Rename: copy sites to new name, delete old
            cats[new_name] = cats.pop(edit_cat)
            msg = _cat_status_html(f"✅ Renamed \"{edit_cat}\" → \"{new_name}\"")
        elif new_name not in cats:
            # New category
            cats[new_name] = [s.strip() for s in (edit_sites or []) if s and s.strip()]
            msg = _cat_status_html(f"✅ Created new category \"{new_name}\" with {len(cats[new_name])} sites")
        else:
            msg = _cat_status_html(f"⚠️ Category \"{new_name}\" already exists", "#e67e22")

        save_categories(cats)
        cat_names = list(cats.keys())
        return (gr.update(choices=cat_names, value=new_name),
                gr.update(choices=cat_names, value=new_name),
                gr.update(value=cats.get(new_name, [])),
                msg)

    def on_cat_delete(cat_name):
        """Delete a category."""
        cats = load_categories()
        if cat_name in cats and len(cats) > 1:
            del cats[cat_name]
            save_categories(cats)
            msg = _cat_status_html(f"✅ Deleted category \"{cat_name}\"")
        elif len(cats) <= 1:
            msg = _cat_status_html(f"⚠️ Cannot delete the last category", "#e74c3c")
        else:
            msg = _cat_status_html(f"⚠️ Category \"{cat_name}\" not found", "#e74c3c")
        cat_names = list(cats.keys())
        first = cat_names[0] if cat_names else ""
        return (gr.update(choices=cat_names, value=first),
                gr.update(choices=cat_names, value=first),
                gr.update(value=cats.get(first, [])),
                gr.update(value=cats.get(first, [])),
                msg)

    def on_cat_save(cat_name, edit_sites):
        """Save edited sites to the category."""
        cats = load_categories()
        sites = [s.strip() for s in (edit_sites or []) if s and s.strip()]
        cats[cat_name] = sites
        save_categories(cats)
        msg = _cat_status_html(f"✅ Saved {len(sites)} sites to \"{cat_name}\"")
        return gr.update(value=sites), msg

    provider_input.change(on_provider_change, inputs=[provider_input], outputs=[provider_url, model_input])
    provider_url.submit(on_provider_url_change, inputs=[provider_url], outputs=[model_input])
    model_timer.tick(poll_models, inputs=[model_input], outputs=[model_input])
    site_category.change(on_category_change, inputs=[site_category, site_category_all], outputs=[research_site])
    site_category_all.change(on_all_categories_toggle, inputs=[site_category_all, site_category], outputs=[research_site])
    research_site.change(on_sites_change, inputs=[research_site])
    cookie_load_btn.click(on_load_cookies, inputs=[cookie_browser, research_site], outputs=[cookies_input, cookie_status])

    # Category management events
    cat_edit_select.change(on_cat_edit_select, inputs=[cat_edit_select], outputs=[cat_edit_sites])
    cat_add_btn.click(on_cat_add_rename, inputs=[cat_edit_select, cat_new_name, cat_edit_sites],
                      outputs=[site_category, cat_edit_select, cat_edit_sites, cat_status])
    cat_del_btn.click(on_cat_delete, inputs=[cat_edit_select],
                      outputs=[site_category, cat_edit_select, cat_edit_sites, research_site, cat_status])
    cat_save_btn.click(on_cat_save, inputs=[cat_edit_select, cat_edit_sites], outputs=[research_site, cat_status])

    def run_research_and_maybe_keynote(prompt, site_url, max_articles, model, cookies, auto_gen, l1, l2, slides, thm, depth):
        """Run research, then optionally auto-generate keynote from top result."""
        global _active_stop_flag
        sf = StopFlag()
        _active_stop_flag = sf
        _progress.clear("research")

        # Save preferences
        is_deep = depth == "Deep"
        save_prefs(max_articles=int(max_articles or 30), auto_keynote=bool(auto_gen),
                   slides=int(slides or 10), lang1=l1, lang2=l2, theme=thm, crawl_depth=depth)

        # 7 outputs: research_log, research_kw, research_out, article_picker, pick_log, pick_file, pick_summary
        no_kn = gr.update()
        last_res_log = ""
        last_kw = ""
        last_out = ""
        last_picker = no_kn
        top_selection = None

        try:
            for res_log, res_kw, res_out, picker_update in run_research(prompt, site_url, int(max_articles), model, cookies, crawl_deep=is_deep, stop_flag=sf):
                sf.check()
                last_res_log = res_log
                last_kw = res_kw
                last_out = res_out
                last_picker = picker_update
                _progress.update(research_log=res_log, research_kw=res_kw, research_out=res_out)
                # Capture the top selection from the final picker update
                if isinstance(picker_update, dict) and picker_update.get("value"):
                    top_selection = picker_update["value"]
                yield res_log, res_kw, res_out, picker_update, no_kn, no_kn, no_kn
        except GeneratorExit:
            last_res_log += "\n\n🛑 Stopped by user"
            _progress.update(research_log=last_res_log)
            _progress.finish()
            yield last_res_log, last_kw, last_out, last_picker, no_kn, no_kn, no_kn
            return
        finally:
            _active_stop_flag = None

        if not auto_gen:
            _progress.finish()
            return

        # Auto-generate from top article
        if not top_selection or "|" not in top_selection:
            last_res_log += "\n\n⚠️ Auto-generate: no article found to generate from"
            _progress.update(research_log=last_res_log)
            _progress.finish()
            yield last_res_log, last_kw, last_out, last_picker, "⚠️ No article to generate from", no_kn, no_kn
            return

        last_pipe = top_selection.rfind("|")
        url = top_selection[last_pipe + 1:].strip()
        title = top_selection[:last_pipe].strip()
        if not url.startswith("http"):
            last_res_log += f"\n\n⚠️ Auto-generate: invalid URL extracted: {url}"
            _progress.update(research_log=last_res_log)
            _progress.finish()
            yield last_res_log, last_kw, last_out, last_picker, f"⚠️ Invalid URL: {url}", no_kn, no_kn
            return

        num = int(slides or 10)

        # Gather top fetchable articles from research results
        extra_urls = []
        top_domain = urlparse(normalize_url(url)).netloc
        domain_counts = {top_domain: 1}
        max_per_domain = 2
        max_extras = 6
        for r in _research_results:
            if r.get("url") == url:
                continue
            r_url = r.get("url", "")
            r_domain = urlparse(normalize_url(r_url)).netloc.removeprefix("www.")
            # Resolve Google News URLs to actual article URLs
            r_url = resolve_google_news_url(r_url)
            r_domain = urlparse(normalize_url(r_url)).netloc.removeprefix("www.")
            # Skip unfetchable
            # Allow up to max_per_domain articles from same domain
            if domain_counts.get(r_domain, 0) >= max_per_domain:
                continue
            extra_urls.append(r_url)
            domain_counts[r_domain] = domain_counts.get(r_domain, 0) + 1
            if len(extra_urls) >= max_extras:
                break

        # Show transition in both logs
        combine_msg = f" + {len(extra_urls)} related" if extra_urls else ""
        last_res_log += f"\n\n🚀 Auto-generating Keynote from: {title[:60]}{combine_msg}..."
        pick_msg = f"🚀 Starting Keynote generation...\n📎 {title[:70]}{combine_msg}\n🔗 {url}"
        _progress.update(research_log=last_res_log, pick_log=pick_msg)
        yield last_res_log, last_kw, last_out, last_picker, pick_msg, no_kn, no_kn

        for kn_log, kn_file, kn_summary in run_pipeline(url, l1, l2, num, thm, True, model, cookies, user_prompt=prompt, extra_urls=extra_urls, stop_flag=sf):
            sf.check()
            _progress.update(pick_log=kn_log)
            yield last_res_log, last_kw, last_out, last_picker, kn_log, kn_file, kn_summary

        _progress.finish()

    research_event = research_btn.click(
        lambda: gr.update(interactive=True),
        outputs=[research_stop],
    ).then(
        run_research_and_maybe_keynote,
        inputs=[research_prompt, research_site, research_max, model_input, cookies_input, auto_keynote, pick_lang1, pick_lang2, pick_slides, pick_theme, crawl_depth],
        outputs=[research_log, research_kw, research_out, article_picker, pick_log, pick_file, pick_summary],
    ).then(
        lambda: gr.update(interactive=False),
        outputs=[research_stop],
    )
    def _trigger_stop():
        global _active_stop_flag
        if _active_stop_flag:
            _active_stop_flag.stop()

    def _stop_research():
        _trigger_stop()
        return gr.update(interactive=False), gr.update(value="🛑 Stopped by user")

    research_stop.click(
        _stop_research,
        outputs=[research_stop, research_log],
        cancels=[research_event],
    )

    def run_from_picker(selection, l1, l2, slides, thm, mdl, cookies, prompt):
        """Extract URL from picker selection and run keynote pipeline."""
        global _active_stop_flag
        sf = StopFlag()
        _active_stop_flag = sf
        _progress.clear("pick")
        try:
            if not selection or "|" not in selection:
                yield "❌ No article selected — run Research first", None, ""
                _progress.finish()
                return
            last_pipe = selection.rfind("|")
            url = selection[last_pipe + 1:].strip()
            if not url.startswith("http"):
                yield f"❌ Invalid URL: {url}", None, ""
                _progress.finish()
                return
            for log, fobj, summary in run_pipeline(url, l1, l2, int(slides), thm, True, mdl, cookies, user_prompt=prompt, stop_flag=sf):
                _progress.update(pick_log=log)
                yield log, fobj, summary
        except GeneratorExit:
            _progress.update(pick_log="🛑 Stopped by user")
            yield "🛑 Stopped by user", None, ""
        except Exception as e:
            yield f"❌ Error: {e}", None, ""
        finally:
            _active_stop_flag = None
            _progress.finish()

    pick_event = generate_from_research.click(
        lambda: gr.update(interactive=True),
        outputs=[pick_stop],
    ).then(
        run_from_picker,
        inputs=[article_picker, pick_lang1, pick_lang2, pick_slides, pick_theme, model_input, cookies_input, research_prompt],
        outputs=[pick_log, pick_file, pick_summary],
    ).then(
        lambda: gr.update(interactive=False),
        outputs=[pick_stop],
    )
    def _stop_pick():
        _trigger_stop()
        return gr.update(interactive=False), gr.update(value="🛑 Stopped by user")

    pick_stop.click(
        _stop_pick,
        outputs=[pick_stop, pick_log],
        cancels=[pick_event],
    )

    def run_pipeline_and_save(url, l1, l2, slides, thm, open_kn, mdl, cookies):
        global _active_stop_flag
        sf = StopFlag()
        _active_stop_flag = sf
        _progress.clear("keynote")
        save_prefs(slides=int(slides or 10), lang1=l1, lang2=l2, theme=thm)
        try:
            for log, fobj, summary in run_pipeline(url, l1, l2, int(slides), thm, open_kn, mdl, cookies, stop_flag=sf):
                _progress.update(keynote_log=log)
                yield log, fobj, summary
        except GeneratorExit:
            _progress.update(keynote_log="🛑 Stopped by user")
            yield "🛑 Stopped by user", None, ""
        finally:
            _active_stop_flag = None
            _progress.finish()

    keynote_event = btn.click(
        lambda: gr.update(interactive=True),
        outputs=[keynote_stop],
    ).then(
        run_pipeline_and_save,
        inputs=[url_input, lang1, lang2, num_slides, theme, open_keynote, model_input, cookies_input],
        outputs=[log_output, file_output, summary_output],
    ).then(
        lambda: gr.update(interactive=False),
        outputs=[keynote_stop],
    )
    def _stop_keynote():
        _trigger_stop()
        return gr.update(interactive=False), gr.update(value="🛑 Stopped by user")

    keynote_stop.click(
        _stop_keynote,
        outputs=[keynote_stop, log_output],
        cancels=[keynote_event],
    )

    # ── Output Files handlers ──
    def _refresh_file_list():
        return gr.update(choices=list_output_filenames(), value=None), ""

    def _resolve_picked(fname):
        if not fname:
            return None
        fpath = os.path.join(OUTPUT_DIR, fname)
        return fpath if os.path.exists(fpath) else None

    def _download_file(fname):
        fpath = _resolve_picked(fname)
        if not fpath:
            return '<span style="color:#e67e22">⚠️ No file selected</span>', gr.update(visible=False, value=None)
        return f'<span style="color:#2ecc71">✅ {fname}</span>', gr.update(visible=True, value=fpath)

    def _open_file(fname):
        fpath = _resolve_picked(fname)
        if not fpath:
            return '<span style="color:#e67e22">⚠️ No file selected</span>'
        subprocess.Popen(["open", fpath])
        return f'<span style="color:#2ecc71">✅ Opened: {fname}</span>'

    def _ask_delete(fname):
        fpath = _resolve_picked(fname)
        if not fpath:
            return '<span style="color:#e67e22">⚠️ No file selected</span>', gr.update(visible=False)
        return f'<span style="color:#e57373;font-size:12px;">Delete {fname}?</span>', gr.update(visible=True)

    def _cancel_delete():
        return '', gr.update(visible=False)

    def _confirm_delete(fname):
        fpath = _resolve_picked(fname)
        if not fpath:
            return gr.update(), '<span style="color:#e67e22">⚠️ No file selected</span>', gr.update(visible=False)
        try:
            os.remove(fpath)
            status = f'<span style="color:#2ecc71">✅ Deleted: {fname}</span>'
        except Exception as e:
            status = f'<span style="color:#e74c3c">❌ {e}</span>'
        return gr.update(choices=list_output_filenames(), value=None), status, gr.update(visible=False)

    output_refresh_btn.click(_refresh_file_list, outputs=[output_file_picker, output_status])
    output_download_btn.click(_download_file, inputs=[output_file_picker], outputs=[output_status, output_download])
    output_open_btn.click(_open_file, inputs=[output_file_picker], outputs=[output_status])
    output_delete_btn.click(_ask_delete, inputs=[output_file_picker], outputs=[output_status, output_confirm_row])
    output_cancel_btn.click(_cancel_delete, outputs=[output_status, output_confirm_row])
    output_confirm_btn.click(_confirm_delete, inputs=[output_file_picker], outputs=[output_file_picker, output_status, output_confirm_row])

    # Auto-refresh after any pipeline completes
    research_event.then(_refresh_file_list, outputs=[output_file_picker, output_status])
    pick_event.then(_refresh_file_list, outputs=[output_file_picker, output_status])
    keynote_event.then(_refresh_file_list, outputs=[output_file_picker, output_status])

if __name__ == "__main__":
    app.launch(
        server_name="0.0.0.0",
        server_port=7860,
        share=False,
        allowed_paths=[OUTPUT_DIR],
        theme=gr.themes.Soft(),
        css=css,
        js=progress_js,
        head=mobile_head,
    )
